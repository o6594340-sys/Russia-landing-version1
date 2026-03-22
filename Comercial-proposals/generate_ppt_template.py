from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "Коммерческое предложение: Япония"
subtitle = slide.placeholders[1]
subtitle.text = "Корпоративное мероприятие / Конференция / Инсентив\nКомпания: [Название клиента]\nДата: [Дата]"

# Slide 2: Concept
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Концепция мероприятия"
content = slide.placeholders[1]
content.text = "Япония — идеальное место для вашего корпоративного события:\n- Культура уважения и дисциплины\n- Современные технологии и традиции\n- Активности, адаптированные под вашу ЦА\n- Максимальная проработка деталей для незабываемого опыта"

# Slide 3: Program Day 1
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Программа: День 1 — Прибытие и знакомство"
content = slide.placeholders[1]
content.text = "Утро: Прибытие в Токио, трансфер в отель\nОбед: Традиционная японская кухня\nВечер: Лекция о японской культуре, адаптированная под вашу команду\nВдохновение: Погружение в атмосферу уважения и инноваций"

# Slide 4: Program Day 2
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Программа: День 2 — Активности и конференции"
content = slide.placeholders[1]
content.text = "Утро: Экскурсия по храмам\nДень: Конференция с лекторами\nВечер: Ужин в ресторане с видом на город\nДетали: Все активности подобраны под стиль клиента"

# Slide 5: Hotels
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Отели"
content = slide.placeholders[1]
content.text = "Отель: [Название], 4*\nРасположение: Центр Токио\nУдобства: Комфортные номера, бизнес-центр\nПитание: Завтраки включены, обеды/ужины по программе"

# Slide 6: Summary
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Итоги и контакты"
content = slide.placeholders[1]
content.text = "Общая стоимость: [Итого из Excel]\nУсловия оплаты: Предоплата 50%\nКонтакты: [Ваши данные]\nГотовы обсудить детали!"

# Save
prs.save("Japan-Proposal-Template.pptx")
print("Шаблон презентации создан: Japan-Proposal-Template.pptx")