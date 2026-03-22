import openpyxl
from pptx import Presentation
import sys

# Function to generate Excel
def generate_excel(num_people, num_days, hotel_level, single_rooms=True):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Смета Япония"

    # Headers
    headers = ["Сервис", "Описание", "Кол-во", "Цена за ед. (USD)", "Итого (USD)"]
    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = openpyxl.styles.Font(bold=True)

    # Data
    hotel_price = 250 if hotel_level == 'premium' else 200
    if single_rooms:
        hotel_price *= 1.5  # Single room surcharge

    data = [
        ["Отель", f"Отель в Токио, {hotel_level}* (одноместное размещение)", num_people, hotel_price, f"=C2*D2"],
        ["Питание", "Завтрак + обед + ужин (японская кухня)", num_people * num_days, 50, f"=C3*D3"],
        ["Экскурсии/Активности", "Культурные экскурсии, лекции (под активных мужчин)", num_people, 100, f"=C4*D4"],
        ["Конференция", "Конференция 3 часа в отеле (1 день)", 1, 2000, f"=C5*D5"],  # Flat fee
        ["Трансфер", "Аэропорт - отель - аэропорт", num_people, 30, f"=C6*D6"],
        ["Итого", "", "", "", "=SUM(E2:E5)"]
    ]

    for row, row_data in enumerate(data, 2):
        for col, value in enumerate(row_data, 1):
            ws.cell(row=row, column=col, value=value)

    wb.save(f'Proposal-{num_people}people-{num_days}days.xlsx')
    print(f"Excel создан: Proposal-{num_people}people-{num_days}days.xlsx")

# Function to generate PPT
def generate_ppt(num_days, conference_day=2, concept_text="Корпоративное мероприятие для активных мужчин: Япония сочетает традиции и современность, идеально для командного духа."):
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Коммерческое предложение: Япония"
    slide.placeholders[1].text = "Корпоративное мероприятие\n100 участников, одноместные номера\n4 дня после 15 мая"

    # Concept slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Концепция"
    slide.placeholders[1].text = concept_text + "\nАктивности адаптированы под мужчин: культурные экскурсии, лекции, без экстремального спорта."

    # Program slides
    activities = [
        "День 1: Прибытие, трансфер в отель, ужин с видом на город.",
        "День 2: Конференция 3 часа в отеле утром, затем экскурсия по Токио (храмы, музеи).",
        "День 3: Активный день: посещение парков, лекция о японской культуре, ужин.",
        "День 4: Свободное время или дополнительная экскурсия, трансфер в аэропорт."
    ]

    for day in range(1, num_days + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"День {day}"
        slide.placeholders[1].text = activities[day-1]

    # Hotels slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Отели"
    slide.placeholders[1].text = "Отель 4*: Одноместные номера, комфорт, бизнес-центр.\nРасположение: Центр Токио."

    # Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Итоги"
    slide.placeholders[1].text = "Общая стоимость: Рассчитана в Excel\nУсловия: Предоплата 50%\nКонтакты: [Ваши данные]"

    prs.save(f'Proposal-{num_days}days-conference.pptx')
    print(f"PPT создан: Proposal-{num_days}days-conference.pptx")

if __name__ == "__main__":
    num_people = int(sys.argv[1]) if len(sys.argv) > 1 else 100
    num_days = int(sys.argv[2]) if len(sys.argv) > 2 else 4
    hotel_level = sys.argv[3] if len(sys.argv) > 3 else 'standard'
    single_rooms = True  # As per brief
    conference_day = 2  # Assume day 2

    generate_excel(num_people, num_days, hotel_level, single_rooms)
    generate_ppt(num_days, conference_day)
    print("КП сгенерировано для 100 чел, 4 дня, одноместные, конференция!")