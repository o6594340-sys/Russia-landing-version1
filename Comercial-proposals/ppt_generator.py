"""
Генератор PPT-презентации для КП.
Дизайн: Washi & Red — тёплая тема, Georgia + Arial, фото через API.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path
import io
import re

from photo_fetcher import fetch_photo, reset_session
from PIL import Image as _PILImage
from claude_generator import detect_season, ROUTES, DEFAULT_ROUTE

BASE_DIR = Path(__file__).parent

W = 13.33   # slide width  (inches)
H = 7.5     # slide height (inches)
M = 0.83    # margin       (inches)


# ── Палитра Washi & Red ──────────────────────────────────────────────────────
class C:
    BEIGE  = RGBColor(0xF7, 0xF4, 0xEF)   # Warm washi beige (content bg)
    INDIGO = RGBColor(0xC8, 0x10, 0x2E)   # Japanese red (accent)
    GOLD   = RGBColor(0xA8, 0x89, 0x5A)   # Matte gold
    DARK   = RGBColor(0x1A, 0x0A, 0x0D)   # Deep dark (title/closing bg)
    INK    = RGBColor(0x1C, 0x1C, 0x1E)   # Body text
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    MIST   = RGBColor(0xE8, 0xE4, 0xDC)   # Light divider stripe
    STONE  = RGBColor(0x76, 0x76, 0x76)   # Secondary/caption text
    GHOST  = RGBColor(0xDE, 0xDA, 0xD2)   # Ghost day numbers


# ── Низкоуровневые утилиты ────────────────────────────────────────────────────

def _bg(slide, color: RGBColor):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, fill: RGBColor):
    shp = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _set_alpha(shape, opacity: float):
    """Устанавливает прозрачность заливки фигуры. opacity: 0.0–1.0."""
    spPr = shape._element.spPr
    solid = spPr.find(qn('a:solidFill'))
    if solid is None:
        return
    for tag in ('a:srgbClr', 'a:sysClr', 'a:schemeClr'):
        clr = solid.find(qn(tag))
        if clr is not None:
            for a in clr.findall(qn('a:alpha')):
                clr.remove(a)
            alpha = etree.SubElement(clr, qn('a:alpha'))
            alpha.set('val', str(int(opacity * 100000)))
            break


def _line(slide, left, top, width, color: RGBColor, thickness_pt=1.0):
    """Тонкая горизонтальная линия."""
    shp = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Pt(thickness_pt)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _txt(slide, text, left, top, width, height,
         size=12, bold=False, italic=False,
         color: RGBColor = C.INK,
         align=PP_ALIGN.LEFT,
         font='Arial'):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _multiline(slide, paragraphs: list, left, top, width, height,
               size=11, color: RGBColor = C.INK,
               bold=False, italic=False, font='Arial',
               line_spacing=1.4, space_before=0):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True

    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if space_before and i > 0:
            p.space_before = Pt(space_before)
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn as _qn
        from lxml import etree as _etree
        # Set line spacing
        pPr = p._p.get_or_add_pPr()
        lnSpc = _etree.SubElement(pPr, _qn('a:lnSpc'))
        spcPct = _etree.SubElement(lnSpc, _qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))

        run = p.add_run()
        run.text = para
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def _cover_crop(img_bytes: bytes, target_w_in: float, target_h_in: float,
                dpi: int = 150) -> bytes:
    """
    Обрезает фото под нужные пропорции без растяжения (как cover в CSS).
    Масштабирует так чтобы заполнить область, лишнее обрезает по центру.
    """
    target_w_px = int(target_w_in * dpi)
    target_h_px = int(target_h_in * dpi)
    target_ratio = target_w_px / target_h_px

    img = _PILImage.open(io.BytesIO(img_bytes)).convert('RGB')
    src_w, src_h = img.size
    src_ratio = src_w / src_h

    if src_ratio > target_ratio:
        # Фото шире нужного — подгоняем по высоте, обрезаем по ширине
        new_h = target_h_px
        new_w = int(src_w * target_h_px / src_h)
    else:
        # Фото выше нужного — подгоняем по ширине, обрезаем по высоте
        new_w = target_w_px
        new_h = int(src_h * target_w_px / src_w)

    img = img.resize((new_w, new_h), _PILImage.LANCZOS)

    left = (new_w - target_w_px) // 2
    top  = (new_h - target_h_px) // 2
    img  = img.crop((left, top, left + target_w_px, top + target_h_px))

    buf = io.BytesIO()
    img.save(buf, format='JPEG', quality=90)
    buf.seek(0)
    return buf.read()


def _photo(slide, img_bytes: bytes, left, top, width, height):
    """Вставляет фото из bytes с правильными пропорциями (cover-crop)."""
    cropped = _cover_crop(img_bytes, width, height)
    stream = io.BytesIO(cropped)
    return slide.shapes.add_picture(
        stream, Inches(left), Inches(top), Inches(width), Inches(height)
    )


def _photo_bg(slide, img_bytes: bytes, overlay_color: RGBColor = C.DARK,
              overlay_opacity: float = 0.62):
    """Полноэкранное фото + полупрозрачный оверлей."""
    _photo(slide, img_bytes, 0, 0, W, H)
    overlay = _rect(slide, 0, 0, W, H, overlay_color)
    _set_alpha(overlay, overlay_opacity)


def _label(slide, text: str, left, top, width=4.0, color: RGBColor = C.INDIGO):
    """Метка верхнего уровня — маленький жирный текст ALL CAPS."""
    _txt(slide, text.upper(), left, top, width, 0.28,
         size=8, bold=True, color=color, font='Arial')


def _logo(slide, right_align=False):
    logo_path = BASE_DIR / 'static' / 'tozai_logo.png'
    if not logo_path.exists():
        return
    logo_h = Inches(0.38)
    logo_w = Inches(1.4)
    left = Inches(W - M - 1.4) if right_align else Inches(M)
    slide.shapes.add_picture(
        str(logo_path), left, Inches(H - 0.62), logo_w, logo_h
    )


def _footer_bar(slide, text='TOZAI TOURS  ·  DMC Japan  ·  tozai-tours.com'):
    """Тонкая нижняя полоска с текстом."""
    _rect(slide, 0, H - 0.28, W, 0.28, C.MIST)
    _txt(slide, text, M, H - 0.27, W - M * 2, 0.24,
         size=7, color=C.STONE, align=PP_ALIGN.CENTER)


# ── Слайд 1: Титул ───────────────────────────────────────────────────────────

def _slide_title(prs, params: dict, credits: list, season_photo: str = ''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Фото фон — сезонное если известен сезон, иначе нейтральное
    cover_query = season_photo if season_photo else 'Japan aerial landscape'
    img, attr = fetch_photo(cover_query, 'Mount Fuji landscape')
    if img:
        _photo_bg(slide, img, C.DARK, overlay_opacity=0.58)
        if attr:
            credits.append(attr)
    else:
        _bg(slide, C.DARK)

    # Золотые горизонтальные линии
    _rect(slide, M, 0.55, W - M * 2, 0.03, C.GOLD)
    _rect(slide, M, H - 0.65, W - M * 2, 0.03, C.GOLD)

    # Метка сверху
    _label(slide, 'Коммерческое предложение', M, 0.22, 5.0, C.GOLD)

    # Главный заголовок
    _txt(slide, 'ЯПОНИЯ  ·  ТОКИО',
         M, 1.3, 10.5, 1.6,
         size=54, bold=True, color=C.WHITE, font='Georgia')

    # Имя компании
    _txt(slide, params['company_name'],
         M, 3.1, 10.5, 0.9,
         size=24, italic=True, color=C.GOLD, font='Georgia')

    # Детали
    twn = params.get('twn', 0)
    sgl = params.get('sgl', 0)
    room_parts = []
    if twn:
        room_parts.append(f'{twn} Twin')
    if sgl:
        room_parts.append(f'{sgl} Single')
    room_str = '  ·  '.join(room_parts) if room_parts else 'номера уточняются'
    conf = '  ·  с конференцией' if params.get('include_conference') else ''
    details = (f"{params['event_type']}{conf}  ·  {params['pax']} чел.  ·  "
               f"{params['days']} дней  ·  отель {params['hotel_level']}  ·  {room_str}")
    if params.get('dates'):
        details += f"  ·  {params['dates']}"

    _txt(slide, details, M, 4.2, 11.5, 0.5,
         size=11, color=RGBColor(0xCC, 0xC9, 0xC2), font='Arial')

    # Логотип и подпись
    _logo(slide, right_align=True)
    _txt(slide, 'Destination Management Company',
         M, H - 0.62, 5.0, 0.38,
         size=8, italic=True, color=C.STONE)


# ── Слайд 2: Почему Япония ────────────────────────────────────────────────────

# Базовые точки — всегда актуальны
_WHY_JAPAN_BASE = [
    'Omotenashi: японский сервис не знает аналогов в мире. Группа чувствует себя гостями, а не туристами.',
    'Страна без случайностей: пунктуальность, безопасность, чистота. Программа работает как часы.',
    'Гастрономия без компромиссов: 230 мишленовских ресторанов в Токио — больше, чем в Париже.',
    'Контраст, который не стирается: 1300 лет истории и передовые технологии в одном городе.',
    'Самый высокий NPS среди инсентив-направлений: участники после Японии становятся её послами.',
]

# Отраслевые точки — заменяют одну базовую при совпадении
_WHY_JAPAN_INDUSTRY = {
    'it':      'Страна, где инженерная культура стала национальной идеей: от Синкансэна до TeamLab — каждый механизм точен до миллиметра.',
    'tech':    'Страна, где инженерная культура стала национальной идеей: от Синкансэна до TeamLab — каждый механизм точен до миллиметра.',
    'фарма':   'Японская концепция здоровья — не отсутствие болезни, а качество жизни. Онсэн, кухня долголетия, ритм, который восстанавливает.',
    'медиц':   'Японская концепция здоровья — не отсутствие болезни, а качество жизни. Онсэн, кухня долголетия, ритм, который восстанавливает.',
    'банк':    'Церемониальность японского делового этикета: каждый жест — знак уважения. Группа вернётся с другим ощущением протокола.',
    'финанс':  'Церемониальность японского делового этикета: каждый жест — знак уважения. Группа вернётся с другим ощущением протокола.',
    'ритейл':  'Японский ритейл — мировой стандарт: Исэтан, Гиндза Сикс, концепция «сотрудник — лицо бренда» берёт начало здесь.',
    'торговл': 'Японский ритейл — мировой стандарт: Исэтан, Гиндза Сикс, концепция «сотрудник — лицо бренда» берёт начало здесь.',
}

# Сезонные точки — добавляются при известном сезоне
_WHY_JAPAN_SEASON = {
    'spring': 'Весна в Японии — событие, которое планируют за год: цветение сакуры длится 10 дней и меняет облик каждого города.',
    'autumn': 'Осень — лучший сезон для групп: момидзи окрашивает парки в алый, температура комфортная, туристов меньше чем весной.',
    'winter': 'Зима открывает другую Японию: идеальный вид на Фудзи, онсэны в снегу, новогодние святилища — тишина после пика сезона.',
    'summer': 'Лето — сезон фестивалей: Бон Одори, фейерверки над реками, зелень буйная. Программа строится вокруг вечерних мероприятий.',
}


def _build_why_japan_points(params: dict, season: str) -> list:
    """Собирает 5 точек под конкретную группу."""
    points = list(_WHY_JAPAN_BASE)  # копия базового списка

    industry = (params.get('industry') or '').lower()

    # Ищем отраслевую точку — если нашли, заменяем 4-ю базовую (контраст истории/техно)
    industry_point = None
    for key, text in _WHY_JAPAN_INDUSTRY.items():
        if key in industry:
            industry_point = text
            break

    if industry_point:
        points[3] = industry_point  # заменяем "контраст истории/техно"

    # Сезонная точка — заменяет последнюю базовую (NPS)
    season_point = _WHY_JAPAN_SEASON.get(season)
    if season_point:
        points[4] = season_point

    return points


def _slide_why_japan(prs, params: dict, season: str, credits: list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    photo_w = W * 0.52
    text_x = photo_w + 0.1

    # Фото слева
    img, attr = fetch_photo('Japan modern traditional contrast', 'Torii gate forest path')
    if img:
        _photo(slide, img, 0, 0, photo_w, H)
        if attr:
            credits.append(attr)
    else:
        _rect(slide, 0, 0, photo_w, H, C.INDIGO)

    # Правая панель — беж
    _bg(slide, C.BEIGE)
    # Перекрываем фото с правой стороны чистым фоном
    _rect(slide, photo_w, 0, W - photo_w, H, C.BEIGE)

    # Тонкая индиго полоска-разделитель
    _rect(slide, photo_w, 0, 0.04, H, C.INDIGO)

    # Метка
    _label(slide, 'Почему Япония', text_x, 1.2, 5.0, C.INDIGO)
    _line(slide, text_x, 1.6, W - text_x - M * 0.5, C.GOLD, 1.5)

    # Заголовок
    _txt(slide, 'Страна, которая\nменяет команды',
         text_x, 1.7, W - text_x - M * 0.5, 1.4,
         size=22, bold=True, color=C.INK, font='Georgia')

    # Пункты
    points = _build_why_japan_points(params, season)
    _multiline(slide, points,
               text_x, 3.25, W - text_x - M * 0.5, 3.5,
               size=10.5, color=C.INK, font='Arial',
               line_spacing=1.55, space_before=5)

    _footer_bar(slide)
    _logo(slide, right_align=True)


# ── Слайд 3: Концепция ────────────────────────────────────────────────────────

def _slide_concept(prs, content: dict, credits: list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C.BEIGE)

    # Полоса сверху
    _rect(slide, 0, 0, W, 1.15, C.INDIGO)

    # Метка и заголовок в полосе
    _label(slide, 'Концепция', M, 0.12, 4.0, C.GOLD)
    title = content.get('concept_title', 'Токио — место силы вашей команды')
    _txt(slide, title, M, 0.38, W - M * 2 - 0.5, 0.72,
         size=20, bold=True, color=C.WHITE, font='Georgia')

    # Декоративный ghost-элемент (кавычка)
    _txt(slide, '«',
         W - 2.8, 0.9, 2.5, 2.5,
         size=160, bold=False, color=C.GHOST, font='Georgia')

    # Текст концепции
    concept_text = content.get('concept_text', '')
    paras = [p.strip() for p in concept_text.split('\n') if p.strip()]
    if not paras:
        paras = [concept_text]

    _multiline(slide, paras,
               M, 1.35, W - M * 2 - 1.5, 5.5,
               size=13, color=C.INK, font='Arial',
               line_spacing=1.65)

    _line(slide, M, 6.75, 3.0, C.GOLD, 1.5)
    _footer_bar(slide)
    _logo(slide, right_align=True)


# ── Слайд: Обзор программы (тайминг) ─────────────────────────────────────────

def _parse_city(text: str) -> str:
    """Извлекает первый город в скобках из строки, напр. '(Токио)'."""
    m = re.search(r'\(([^)]{2,20})\)', text)
    return m.group(1) if m else ''


def _parse_time_entry(line: str):
    """Разбирает строку '10:00 — Активность' → (time_str, activity_str)."""
    line = re.sub(r'\s*\([^)]{2,20}\)', '', line).strip()
    m = re.match(r'^(\d{1,2}:\d{2}(?:[–—\-]\d{1,2}:\d{2})?)(.*)', line)
    if m:
        time_str = m.group(1).strip()
        act_str  = re.sub(r'^[\s—–\-]+', '', m.group(2)).strip()
        return time_str, act_str
    return '', line


def _overview_collect_entries(day_data: dict) -> list:
    """Собирает все тайминг-строки из morning + afternoon + evening."""
    entries = []
    for key in ('morning', 'afternoon', 'evening'):
        raw = day_data.get(key, '').strip()
        if not raw:
            continue
        for line in raw.split('\n'):
            line = line.strip()
            if line:
                t, a = _parse_time_entry(line)
                entries.append((t, a))
    return entries


def _overview_entry_box(slide, time_str: str, activity: str,
                        left, top, width, height, font_size=9.0):
    """Одна строка тайминга: ВРЕМЯ (красный) + текст (тёмный)."""
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _etree

    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = False

    p = tf.paragraphs[0]
    pPr = p._p.get_or_add_pPr()
    lnSpc = _etree.SubElement(pPr, _qn('a:lnSpc'))
    spcPct = _etree.SubElement(lnSpc, _qn('a:spcPct'))
    spcPct.set('val', '100000')

    if time_str:
        r_t = p.add_run()
        r_t.text = time_str
        r_t.font.name = 'Arial'
        r_t.font.size = Pt(font_size)
        r_t.font.bold = True
        r_t.font.color.rgb = C.INDIGO

        r_sp = p.add_run()
        r_sp.text = '  '
        r_sp.font.name = 'Arial'
        r_sp.font.size = Pt(font_size)

    r_a = p.add_run()
    r_a.text = activity
    r_a.font.name = 'Arial'
    r_a.font.size = Pt(font_size)
    r_a.font.color.rgb = C.INK


def _slide_overview(prs, params: dict, content: dict):
    """Обзорный слайд: сплошной тайминг по дням (время + активность)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C.BEIGE)

    days = content.get('days', [])
    n    = max(len(days), 1)

    # ── Шапка ─────────────────────────────────────────────────────────────────
    HDR_H = 0.82
    _rect(slide, 0, 0, W, HDR_H, C.DARK)
    _rect(slide, M, HDR_H - 0.03, W - M * 2, 0.03, C.INDIGO)
    _label(slide, 'Обзор программы', M, 0.10, 6.0, C.GOLD)

    n_word = 'день' if n == 1 else ('дня' if n in (2, 3, 4) else 'дней')
    route_key = params.get('route', DEFAULT_ROUTE)
    route_cities = ROUTES.get(route_key, ROUTES[DEFAULT_ROUTE])['cities']
    cities_str = ' → '.join(route_cities)
    route_label = f'{cities_str}  ·  {n} {n_word}'
    if params.get('dates'):
        route_label += f'  ·  {params["dates"]}'
    _txt(slide, route_label, M, 0.32, W - M * 2 - 0.3, 0.46,
         size=16, bold=True, color=C.WHITE, font='Georgia')

    # ── Параметры сетки ────────────────────────────────────────────────────────
    GRID_TOP  = HDR_H + 0.06
    GRID_BOT  = H - 0.30
    GRID_H    = GRID_BOT - GRID_TOP
    DAY_HDR_H = 0.46
    COL_GAP   = 0.07
    avail_w   = W - M * 2
    col_w     = avail_w / n

    # Сначала собираем все строки чтобы знать максимум
    all_entries = [_overview_collect_entries(d) for d in days]
    max_e  = max((len(e) for e in all_entries), default=7)
    max_e  = max(max_e, 5)   # не меньше 5 строк, чтобы не было огромных ячеек
    avail_h = GRID_H - DAY_HDR_H - 0.20
    entry_h = min(avail_h / max_e, 0.62)   # не выше 0.62" на строку
    font_sz = 9.0  # минимум 9pt везде

    for di, day_data in enumerate(days):
        x  = M + di * col_w
        cx = x + COL_GAP / 2
        cw = col_w - COL_GAP

        # ── Заголовок колонки ─────────────────────────────────────────────
        day_num   = day_data.get('day_num', di + 1)
        raw_title = day_data.get('title', f'День {day_num}')
        short_t   = re.sub(r'^День\s*\d+\s*[—–-]\s*', '', raw_title).strip()
        city      = (_parse_city(day_data.get('morning', ''))
                     or _parse_city(day_data.get('afternoon', ''))
                     or '')

        hdr_fill  = C.INDIGO if di % 2 == 0 else C.DARK
        _rect(slide, cx, GRID_TOP, cw, DAY_HDR_H, hdr_fill)

        _txt(slide, f'ДЕНЬ {day_num}',
             cx + 0.10, GRID_TOP + 0.04, cw - 0.18, 0.18,
             size=7, bold=True, color=C.GOLD, font='Arial')

        max_ch      = max(12, int(cw * 11))
        disp_title  = (short_t[:max_ch - 1] + '…') if len(short_t) > max_ch else short_t
        title_y     = GRID_TOP + 0.22
        _txt(slide, disp_title, cx + 0.10, title_y, cw - 0.18, 0.22,
             size=8.5, bold=True, color=C.WHITE, font='Arial')

        if city:
            _txt(slide, city.upper(),
                 cx + 0.10, GRID_TOP + 0.30, cw - 0.18, 0.16,
                 size=6.5, bold=True,
                 color=RGBColor(0xCC, 0xC9, 0xC2), font='Arial')

        # Разделитель между колонками
        if di > 0:
            _rect(slide, x, GRID_TOP, 0.014, GRID_H, C.STONE)

        # Фон контентной области
        _rect(slide, cx, GRID_TOP + DAY_HDR_H, cw,
              GRID_H - DAY_HDR_H, C.BEIGE)

        # ── Тайминг-строки ────────────────────────────────────────────────
        entries   = all_entries[di]
        content_y = GRID_TOP + DAY_HDR_H + 0.10

        # Определяем где кончается morning и afternoon для разделителей
        def _section_len(key):
            raw = day_data.get(key, '').strip()
            return len([l for l in raw.split('\n') if l.strip()]) if raw else 0

        morning_len   = _section_len('morning')
        afternoon_len = _section_len('afternoon')

        for ei, (time_str, activity) in enumerate(entries):
            # Тонкий разделитель перед началом午后/вечера
            if ei == morning_len and morning_len > 0 and ei > 0:
                sep_y = content_y + ei * entry_h - 0.04
                _rect(slide, cx + 0.08, sep_y, cw - 0.16, 0.012, C.GOLD)
            elif ei == morning_len + afternoon_len and afternoon_len > 0 and ei > 0:
                sep_y = content_y + ei * entry_h - 0.04
                _rect(slide, cx + 0.08, sep_y, cw - 0.16, 0.012, C.STONE)

            ey = content_y + ei * entry_h

            # Чередующийся фон строки
            row_fill = C.BEIGE if ei % 2 == 0 else C.MIST
            _rect(slide, cx, ey, cw, entry_h - 0.02, row_fill)

            # Контент строки — обрезка по реальной ширине колонки
            max_a = max(8, int(cw * 9))
            act   = (activity[:max_a - 1] + '…') if len(activity) > max_a else activity
            _overview_entry_box(slide, time_str, act,
                                cx + 0.10, ey + 0.03,
                                cw - 0.18, entry_h - 0.06,
                                font_size=font_sz)

    # ── Нижняя линия + футер ──────────────────────────────────────────────────
    _line(slide, M, GRID_BOT - 0.01, W - M * 2, C.GOLD, 1.0)
    _footer_bar(slide)
    _logo(slide, right_align=True)


# ── Слайд: День программы ─────────────────────────────────────────────────────

def _build_day_photo_query(day_data: dict, season: str) -> tuple:
    """Строит поисковый запрос фото по содержимому дня, а не по номеру."""
    full_text = ' '.join([
        day_data.get('title', ''),
        day_data.get('morning', ''),
        day_data.get('afternoon', ''),
        day_data.get('evening', ''),
    ]).lower()

    season_sfx = {
        'spring': 'spring sakura',
        'summer': 'summer green',
        'autumn': 'autumn momiji',
        'winter': 'winter snow',
    }.get(season, '')

    # Конкретные места — наивысший приоритет
    if 'teamlab' in full_text:
        return 'TeamLab digital art installation Tokyo', 'TeamLab Planets light art'
    if 'синкансэн' in full_text or 'shinkansen' in full_text:
        return 'Shinkansen bullet train Mount Fuji', 'bullet train Japan landscape'
    if 'тодайдзи' in full_text or 'todaiji' in full_text:
        return 'Todaiji Great Buddha Nara', 'Nara deer park temple'
    if 'кинкакудзи' in full_text or 'kinkakuji' in full_text or 'золотой павильон' in full_text:
        return 'Kinkakuji golden pavilion Kyoto', 'Kinkakuji reflection pond'
    if 'рёандзи' in full_text or 'ryoanji' in full_text:
        return 'Ryoanji zen garden Kyoto', 'Japanese rock garden'
    if 'фусими инари' in full_text or 'fushimi inari' in full_text:
        return 'Fushimi Inari torii gates Kyoto', 'Fushimi Inari red gates path'
    if 'рёкан' in full_text or 'онсэн' in full_text or 'onsen' in full_text:
        return f'Japanese ryokan onsen {season_sfx}'.strip(), 'ryokan tatami room Mount Fuji'
    if 'цукидзи' in full_text or 'tsukiji' in full_text or 'рыбный рынок' in full_text:
        return 'Tsukiji fish market Tokyo morning', 'Japanese fish market seafood'
    if 'асакуса' in full_text or 'сэнсодзи' in full_text or 'asakusa' in full_text:
        return 'Asakusa Sensoji temple Tokyo', 'Nakamise shopping street Asakusa'
    if 'хаконэ' in full_text or 'романсукар' in full_text or 'romancecar' in full_text:
        return 'Hakone Mount Fuji reflection lake', 'Hakone open air museum'
    if 'сибуя' in full_text or 'shibuya' in full_text:
        return 'Shibuya crossing aerial Tokyo', 'Shibuya night Tokyo street'
    if 'синдзюку' in full_text or 'shinjuku' in full_text:
        return 'Shinjuku Tokyo night neon', 'Shinjuku Gyoen garden Tokyo'

    # Определяем город из тегов вида (Киото), (Токио) ...
    city_m = re.search(r'\(([^)]{2,15})\)', full_text)
    city = city_m.group(1) if city_m else ''

    if 'киото' in city or 'нара' in city:
        return f'Kyoto temple {season_sfx}'.strip(), 'Kyoto traditional street Japan'
    if 'нара' in full_text:
        return 'Nara deer park Japan', 'Nara Todaiji ancient temple'
    if 'осака' in city or 'осака' in full_text:
        return 'Osaka castle Japan', 'Dotonbori Osaka night canal'
    if 'хаконэ' in city:
        return 'Hakone Mount Fuji Japan', 'Hakone lake reflection'

    # Fallback — Токио с учётом сезона
    return f'Tokyo {season_sfx}'.strip() or 'Tokyo cityscape Japan', 'Tokyo street Japan'


def _day_fetch_photo(day_data: dict, season: str, credits: list):
    """Загружает фото для слайда дня по содержимому дня."""
    q1, q2 = _build_day_photo_query(day_data, season)
    img, attr = fetch_photo(q1, q2)
    if attr:
        credits.append(attr)
    return img


def _day_schedule_cols(slide, day_data: dict,
                       y_badge: float, col_w: float, text_h: float):
    """Три колонки расписания (утро / день / вечер) — общий блок.
    Время выделяется цветом секции, активность — светлым."""
    sections = [
        ('Утро',  day_data.get('morning', ''),   C.INDIGO,                    C.WHITE,  C.INDIGO),
        ('День',  day_data.get('afternoon', ''), C.GOLD,                      C.INK,   C.GOLD),
        ('Вечер', day_data.get('evening', ''),   RGBColor(0x3A, 0x3A, 0x3A), C.WHITE, RGBColor(0xB0, 0x9A, 0x6A)),
    ]
    ACT_COLOR = RGBColor(0xE0, 0xDD, 0xD8)
    FONT_SZ   = 9.5
    TIME_W    = 1.0   # фиксированная ширина поля времени

    for i, (label_text, text, badge_fill, badge_text_color, time_color) in enumerate(sections):
        x = M + i * (col_w + 0.1)
        if i > 0:
            _rect(slide, x - 0.06, y_badge, 0.02, text_h + 0.34,
                  RGBColor(0x55, 0x55, 0x55))
        _rect(slide, x, y_badge, 1.05, 0.27, badge_fill)
        _txt(slide, label_text.upper(),
             x, y_badge, 1.05, 0.27,
             size=7.5, bold=True, color=badge_text_color,
             align=PP_ALIGN.CENTER, font='Arial')
        if not text:
            continue

        lines = [l.strip() for l in text.split('\n') if l.strip()]
        lines = [re.sub(r'\s*\([^)]{2,20}\)', '', l).strip() for l in lines]
        n = max(len(lines), 1)
        entry_h = min(text_h / n, 0.62)

        for ei, line in enumerate(lines):
            t, act = _parse_time_entry(line)
            ey = y_badge + 0.34 + ei * entry_h

            if t:
                # Время — цвет секции, жирный
                _txt(slide, t,
                     x, ey, TIME_W, entry_h,
                     size=FONT_SZ, bold=True, color=time_color, font='Arial')
                act_x = x + TIME_W + 0.04
                act_w = col_w - TIME_W - 0.08
            else:
                act_x = x
                act_w = col_w - 0.05

            # Обрезаем длинный текст чтобы не вылезал за границу колонки
            max_ch = max(10, int(act_w * 11))
            act_disp = (act[:max_ch - 1] + '…') if len(act) > max_ch else act

            _txt(slide, act_disp,
                 act_x, ey, act_w, entry_h,
                 size=FONT_SZ, color=ACT_COLOR, font='Arial')


# ── Layout A: полный экран — шапка сверху, расписание снизу ──────────────────

def _slide_day_a(prs, day_data: dict, credits: list, season: str = 'unknown'):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    day_num = day_data.get('day_num', 1)

    img = _day_fetch_photo(day_data, season, credits)
    if img:
        _photo(slide, img, 0, 0, W, H)
    else:
        _bg(slide, C.DARK)

    top = _rect(slide, 0, 0, W, 1.55, C.DARK)
    _set_alpha(top, 0.80)
    bot = _rect(slide, 0, H - 3.1, W, 3.1, C.DARK)
    _set_alpha(bot, 0.88)

    _rect(slide, M, 1.55, W - M * 2, 0.03, C.INDIGO)

    _label(slide, 'Программа', M, 0.14, 4.0, C.GOLD)
    title = day_data.get('title', f'День {day_num}')
    _txt(slide, title, M, 0.42, W - M * 2 - 1.5, 0.95,
         size=22, bold=True, color=C.WHITE, font='Georgia')

    _txt(slide, str(day_num).zfill(2),
         W - 3.0, 0.55, 2.6, 2.6,
         size=120, bold=True, color=C.GHOST, font='Georgia')

    description = day_data.get('description', '')
    if description:
        first_sent = re.split(r'(?<=[.!?])\s', description)[0].strip()
        desc_bg = _rect(slide, M - 0.1, 1.72, W - M * 2 + 0.2, 0.75, C.DARK)
        _set_alpha(desc_bg, 0.40)
        _txt(slide, first_sent,
             M, 1.80, W - M * 2 - 0.3, 0.65,
             size=12, italic=True, color=C.WHITE, font='Georgia')

    col_w = (W - M * 2 - 0.2) / 3
    _day_schedule_cols(slide, day_data, y_badge=H - 2.95,
                       col_w=col_w, text_h=2.5)
    _logo(slide, right_align=True)


# ── Layout B: без фото — чистая инфографика, три колонки на светлом фоне ──────

def _slide_day_b(prs, day_data: dict, credits: list, season: str = 'unknown'):
    """Текстовый слайд без фото: три колонки тайминга на бежевом фоне."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    day_num = day_data.get('day_num', 1)

    _bg(slide, C.BEIGE)

    # ── Шапка ────────────────────────────────────────────────────────────────
    HDR_H = 1.30
    _rect(slide, 0, 0, W, HDR_H, C.DARK)
    _rect(slide, M, HDR_H, W - M * 2, 0.03, C.INDIGO)

    _label(slide, 'Программа', M, 0.14, 4.0, C.GOLD)
    title = day_data.get('title', f'День {day_num}')
    _txt(slide, title, M, 0.40, W - M * 2 - 2.8, 0.80,
         size=24, bold=True, color=C.WHITE, font='Georgia')

    description = day_data.get('description', '')
    if description:
        first_sent = re.split(r'(?<=[.!?])\s', description)[0].strip()
        _txt(slide, first_sent,
             M, 0.95, W - M * 2 - 2.8, 0.42,
             size=11, italic=True,
             color=RGBColor(0xCC, 0xC9, 0xC2), font='Georgia')

    # Ghost номер справа в шапке
    _txt(slide, str(day_num).zfill(2),
         W - 2.65, -0.15, 2.5, 1.55,
         size=120, bold=True, color=RGBColor(0x30, 0x22, 0x24), font='Georgia')

    # ── Три колонки тайминга ─────────────────────────────────────────────────
    SCHED_TOP  = HDR_H + 0.20
    SCHED_BOT  = H - 0.35
    SCHED_H    = SCHED_BOT - SCHED_TOP
    COL_GAP    = 0.28
    col_w      = (W - M * 2 - COL_GAP * 2) / 3

    SECTIONS = [
        ('Утро',   'morning',   C.INDIGO, C.WHITE,                        C.INDIGO),
        ('День',   'afternoon', C.GOLD,   C.INK,                          C.GOLD),
        ('Вечер',  'evening',   C.INK,    RGBColor(0xF0, 0xED, 0xE8),    RGBColor(0xB0, 0x9A, 0x6A)),
    ]

    for i, (label_text, key, hdr_fill, hdr_text, time_color) in enumerate(SECTIONS):
        cx = M + i * (col_w + COL_GAP)

        # Цветная шапка колонки
        _rect(slide, cx, SCHED_TOP, col_w, 0.34, hdr_fill)
        _txt(slide, label_text.upper(),
             cx + 0.12, SCHED_TOP, col_w - 0.18, 0.34,
             size=9, bold=True, color=hdr_text, font='Arial')

        # Левая цветная полоска-акцент
        _rect(slide, cx, SCHED_TOP + 0.34, 0.04, SCHED_H - 0.34, hdr_fill)

        # Фон колонки
        _rect(slide, cx + 0.04, SCHED_TOP + 0.34, col_w - 0.04,
              SCHED_H - 0.34, C.BEIGE)

        # Строки тайминга
        raw = day_data.get(key, '').strip()
        if not raw:
            continue
        lines = [l.strip() for l in raw.split('\n') if l.strip()]
        lines = [re.sub(r'\s*\([^)]{2,20}\)', '', l).strip() for l in lines]

        # Динамическая высота строки
        entry_h = min((SCHED_H - 0.44) / max(len(lines), 1), 0.68)
        font_sz = 9.5 if entry_h >= 0.44 else 9.0  # минимум 9pt

        TIME_W = 1.0  # фиксированная ширина поля времени (дюймы)

        for ei, line in enumerate(lines):
            t, act = _parse_time_entry(line)
            ey = SCHED_TOP + 0.44 + ei * entry_h

            # Чередующийся фон
            if ei % 2 == 1:
                _rect(slide, cx + 0.04, ey - 0.03, col_w - 0.04,
                      entry_h, C.MIST)

            # Время — цветной акцент, фиксированная ширина
            if t:
                _txt(slide, t,
                     cx + 0.14, ey, TIME_W, entry_h,
                     size=font_sz, bold=True, color=time_color, font='Arial')
                act_x = cx + 0.14 + TIME_W + 0.05
                act_w = col_w - 0.14 - TIME_W - 0.14
            else:
                act_x = cx + 0.14
                act_w = col_w - 0.22

            max_ch = max(12, int(act_w * 12))
            act_disp = (act[:max_ch - 1] + '…') if len(act) > max_ch else act
            _txt(slide, act_disp,
                 act_x, ey, act_w, entry_h,
                 size=font_sz, color=C.INK, font='Arial')

    _line(slide, M, SCHED_BOT + 0.01, W - M * 2, C.GOLD, 1.0)
    _footer_bar(slide)
    _logo(slide, right_align=True)


# ── Layout C: фото сверху, чистая тёмная панель снизу ────────────────────────

def _slide_day_c(prs, day_data: dict, credits: list, season: str = 'unknown'):
    """Горизонтальный сплит: фото верхние 55%, тёмная панель с расписанием снизу."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    day_num = day_data.get('day_num', 1)

    PHOTO_H = H * 0.54     # ~4.05"
    PANEL_Y = PHOTO_H

    img = _day_fetch_photo(day_data, season, credits)
    if img:
        _photo(slide, img, 0, 0, W, PHOTO_H)
    else:
        _rect(slide, 0, 0, W, PHOTO_H, C.DARK)

    # Тёмная панель снизу
    _bg(slide, C.DARK)
    _rect(slide, 0, PANEL_Y, W, H - PANEL_Y, RGBColor(0x12, 0x08, 0x09))

    # Красная линия на стыке
    _rect(slide, 0, PANEL_Y, W, 0.04, C.INDIGO)

    # Заголовок дня — плавает поверх фото у нижнего края
    title_y = PHOTO_H - 1.30
    title_bg = _rect(slide, 0, title_y, W, 1.26, C.DARK)
    _set_alpha(title_bg, 0.72)

    _label(slide, 'Программа', M, title_y + 0.10, 4.0, C.GOLD)
    title = day_data.get('title', f'День {day_num}')
    _txt(slide, title, M, title_y + 0.34, W - M * 2 - 1.5, 0.85,
         size=22, bold=True, color=C.WHITE, font='Georgia')

    # Ghost номер — правый верхний угол фото
    _txt(slide, str(day_num).zfill(2),
         W - 3.0, 0.15, 2.8, 2.6,
         size=120, bold=True, color=C.GHOST, font='Georgia')

    # Описание — в тёмной панели, только первое предложение (места мало)
    description = day_data.get('description', '')
    if description:
        first_sent = re.split(r'(?<=[.!?])\s', description)[0].strip()
        _txt(slide, first_sent,
             M, PANEL_Y + 0.14, W - M * 2 - 0.5, 0.55,
             size=11, italic=True, color=RGBColor(0xCC, 0xC9, 0xC2), font='Georgia')

    # Три колонки расписания в нижней панели
    col_w  = (W - M * 2 - 0.2) / 3
    sched_y = PANEL_Y + (0.80 if description else 0.22)
    text_h  = H - sched_y - 0.55
    _day_schedule_cols(slide, day_data, y_badge=sched_y,
                       col_w=col_w, text_h=text_h)

    _logo(slide, right_align=True)


# ── Диспетчер: выбирает layout по номеру дня ─────────────────────────────────

_DAY_LAYOUTS = [_slide_day_a, _slide_day_b, _slide_day_c]


def _slide_day(prs, day_data: dict, credits: list, season: str = 'unknown'):
    day_num = day_data.get('day_num', 1)
    _DAY_LAYOUTS[(day_num - 1) % len(_DAY_LAYOUTS)](prs, day_data, credits, season)


# ── Слайд: Хайлайт — три варианта вёрстки ────────────────────────────────────

def _hl_label_and_text(slide, hl: dict, text_x: float, panel_w: float, logo_right: bool):
    """Общий блок: категория + заголовок + линия + описание."""
    category   = hl.get('category', '')
    day_ref    = hl.get('day_ref', '')
    label_text = f'{category}  ·  {day_ref}' if category and day_ref else category or day_ref
    if label_text:
        _label(slide, label_text, text_x, 1.0, panel_w - M, C.GOLD)

    title = hl.get('title', '')
    _txt(slide, title, text_x, 1.35, panel_w - M * 1.2, 1.6,
         size=30, bold=True, color=C.WHITE, font='Georgia')

    _line(slide, text_x, 3.1, panel_w - M, C.INDIGO, 1.5)

    description = hl.get('description', '')
    if description:
        paras = [p.strip() for p in description.split('\n') if p.strip()]
        _multiline(slide, paras, text_x, 3.25, panel_w - M * 1.2, 3.6,
                   size=11.5, color=RGBColor(0xE0, 0xDD, 0xD8),
                   font='Arial', line_spacing=1.65)

    _logo(slide, right_align=logo_right)


def _hl_layout_a(prs, hl: dict, img: bytes, credits: list):
    """Вариант А: фото справа, текст слева."""
    slide   = prs.slides.add_slide(prs.slide_layouts[6])
    panel_w = W * 0.52

    if img:
        _photo(slide, img, panel_w, 0, W - panel_w, H)
    else:
        _rect(slide, panel_w, 0, W - panel_w, H, C.DARK)

    _bg(slide, C.DARK)
    _rect(slide, 0, 0, panel_w, H, C.DARK)
    panel = _rect(slide, 0, 0, panel_w, H, RGBColor(0x12, 0x08, 0x09))
    _set_alpha(panel, 0.92)

    right_vignette = _rect(slide, panel_w, 0, W - panel_w, H, C.DARK)
    _set_alpha(right_vignette, 0.15)

    _rect(slide, panel_w - 0.04, 0, 0.04, H, C.INDIGO)
    _hl_label_and_text(slide, hl, M, panel_w, logo_right=False)


def _hl_layout_b(prs, hl: dict, img: bytes, credits: list):
    """Вариант Б: фото слева, текст справа."""
    slide   = prs.slides.add_slide(prs.slide_layouts[6])
    photo_w = W * 0.48
    text_x  = photo_w + 0.14
    panel_w = W - text_x - 0.1

    if img:
        _photo(slide, img, 0, 0, photo_w, H)
    else:
        _rect(slide, 0, 0, photo_w, H, C.DARK)

    _bg(slide, RGBColor(0x12, 0x08, 0x09))
    _rect(slide, photo_w, 0, W - photo_w, H, RGBColor(0x12, 0x08, 0x09))

    left_vignette = _rect(slide, 0, 0, photo_w, H, C.DARK)
    _set_alpha(left_vignette, 0.15)

    _rect(slide, photo_w, 0, 0.04, H, C.INDIGO)
    _hl_label_and_text(slide, hl, text_x, panel_w, logo_right=True)


def _hl_layout_c(prs, hl: dict, img: bytes, credits: list):
    """Вариант В: фото на весь слайд, тёмная плашка снизу с текстом."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if img:
        _photo(slide, img, 0, 0, W, H)
    else:
        _bg(slide, C.DARK)

    # Тёмная плашка снизу
    plashka_h = 3.0
    plashka_y = H - plashka_h
    plashka = _rect(slide, 0, plashka_y, W, plashka_h, C.DARK)
    _set_alpha(plashka, 0.88)

    # Красная линия сверху плашки
    _rect(slide, M, plashka_y + 0.01, W - M * 2, 0.03, C.INDIGO)

    # Категория + день
    category   = hl.get('category', '')
    day_ref    = hl.get('day_ref', '')
    label_text = f'{category}  ·  {day_ref}' if category and day_ref else category or day_ref
    if label_text:
        _label(slide, label_text, M, plashka_y + 0.18, W - M * 2, C.GOLD)

    # Заголовок крупно
    title = hl.get('title', '')
    _txt(slide, title, M, plashka_y + 0.45, W * 0.55, 1.0,
         size=28, bold=True, color=C.WHITE, font='Georgia')

    # Описание — две колонки на широкой плашке
    description = hl.get('description', '')
    if description:
        paras = [p.strip() for p in description.split('\n') if p.strip()]
        mid = len(paras) // 2 + len(paras) % 2
        col_w = W / 2 - M - 0.1
        _multiline(slide, paras[:mid], M, plashka_y + 1.55, col_w, 1.3,
                   size=10.5, color=RGBColor(0xE0, 0xDD, 0xD8),
                   font='Arial', line_spacing=1.5)
        if paras[mid:]:
            _multiline(slide, paras[mid:], W / 2 + 0.1, plashka_y + 1.55, col_w, 1.3,
                       size=10.5, color=RGBColor(0xE0, 0xDD, 0xD8),
                       font='Arial', line_spacing=1.5)

    _logo(slide, right_align=True)


_HL_LAYOUTS = [_hl_layout_a, _hl_layout_b, _hl_layout_c]


def _slide_highlight(prs, hl: dict, credits: list, index: int = 0):
    """Выбирает вариант вёрстки по индексу (A → Б → В → A...)."""
    query = hl.get('photo_query', hl.get('title', 'Japan'))
    img, attr = fetch_photo(query, query)
    if attr:
        credits.append(attr)

    layout_fn = _HL_LAYOUTS[index % len(_HL_LAYOUTS)]
    layout_fn(prs, hl, img, credits)


# ── Слайд: Размещение ─────────────────────────────────────────────────────────

_HOTEL_PHOTO_QUERIES = [
    ('Japanese luxury hotel lobby interior', 'Tokyo hotel room modern'),
    ('Tokyo hotel garden view room', 'Japan hotel panoramic window'),
]


def _slide_hotel_card(prs, params: dict, content: dict, credits: list,
                      variant_label: str, hotel_name: str,
                      show_closing: bool = False, photo_index: int = 0):
    """Один слайд отеля. variant_label = 'ОПЦИЯ А — Prince Park Tower' и т.п."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    photo_w = W * 0.48
    text_x  = photo_w + 0.12

    q1, q2 = _HOTEL_PHOTO_QUERIES[photo_index % len(_HOTEL_PHOTO_QUERIES)]
    img, attr = fetch_photo(q1, q2)
    if img:
        _photo(slide, img, 0, 0, photo_w, H)
        if attr:
            credits.append(attr)
    else:
        _rect(slide, 0, 0, photo_w, H, C.INDIGO)

    _bg(slide, C.BEIGE)
    _rect(slide, photo_w, 0, W - photo_w, H, C.BEIGE)
    _rect(slide, photo_w, 0, 0.04, H, C.INDIGO)

    # Метка «ОПЦИЯ А / Б» + название отеля
    _label(slide, variant_label, text_x, 0.72, W - text_x - 0.5, C.INDIGO)
    _line(slide, text_x, 1.10, W - text_x - M * 0.5, C.GOLD, 1.5)
    _txt(slide, hotel_name,
         text_x, 1.18, W - text_x - 0.5, 1.0,
         size=22, bold=True, color=C.INK, font='Georgia')

    # Описание отеля (из enrich_texts; только на первом слайде)
    hotel_desc = content.get('hotel_description', '') if photo_index == 0 else ''
    paras = [p.strip() for p in hotel_desc.split('\n') if p.strip()] if hotel_desc else []
    if paras:
        _multiline(slide, paras,
                   text_x, 2.35, W - text_x - 0.5, 2.0,
                   size=11.5, color=C.INK, font='Arial', line_spacing=1.55)

    # Параметры
    twn    = params.get('twn', 0)
    sgl    = params.get('sgl', 0)
    nights = params['days'] - 1
    room_parts = []
    if twn: room_parts.append(f'{twn} Twin-номеров')
    if sgl: room_parts.append(f'{sgl} Single-номеров')
    room_str = ' + '.join(room_parts) if room_parts else 'уточняется'

    params_lines = [
        f'Состав номеров  —  {room_str}',
        f'Количество ночей  —  {nights}',
        f'Гостей  —  {params["pax"]} человек',
        f'Уровень  —  {params["hotel_level"]}',
    ]
    params_top = 4.5 if paras else 2.35
    _multiline(slide, params_lines,
               text_x, params_top, W - text_x - 0.5, 1.8,
               size=10, color=C.STONE, font='Arial',
               line_spacing=1.6, space_before=4)

    # Closing note — только на последнем слайде размещения
    if show_closing:
        closing = content.get('closing_note', '')
        if closing:
            _rect(slide, text_x - 0.15, 6.3, W - text_x, 0.85, C.MIST)
            _txt(slide, closing, text_x, 6.35, W - text_x - 0.3, 0.75,
                 size=10, italic=True, color=C.INK, font='Georgia')

    _footer_bar(slide)
    _logo(slide, right_align=True)


def _slide_hotels(prs, params: dict, content: dict, credits: list):
    """Один или два слайда размещения в зависимости от наличия второго отеля."""
    hotel_a = (params.get('hotel_a_name') or 'Отель').strip()
    hotel_b = (params.get('hotel_b_name') or '').strip()

    has_b = bool(hotel_b and params.get('hotel_b_rate', 0))

    _slide_hotel_card(prs, params, content, credits,
                      variant_label=f'Опция А  —  Размещение',
                      hotel_name=hotel_a,
                      show_closing=not has_b,
                      photo_index=0)

    if has_b:
        _slide_hotel_card(prs, params, content, credits,
                          variant_label=f'Опция Б  —  Альтернативное размещение',
                          hotel_name=hotel_b,
                          show_closing=True,
                          photo_index=1)


# ── Слайд: Финал / Контакты ───────────────────────────────────────────────────

# Контакты подставляются при первом вызове из Tozai ОБЩАЯ.pptx,
# или используется значение по умолчанию.
_DEFAULT_CONTACTS = {
    'name': 'Mr. Hiroshi Suhara',
    'title': 'MICE Director · Tozai Tours DMC Japan',
    'email': 'suhara@tozaitours.com',
    'phone': '+81 3 6403 4816',
    'web': 'tozai-tours.com',
}

_CONTACTS: dict | None = None


def set_contacts(contacts: dict):
    """Вызывается из app.py после парсинга контактов из ОБЩАЯ.pptx."""
    global _CONTACTS
    _CONTACTS = contacts


def _get_contacts() -> dict:
    return _CONTACTS or _DEFAULT_CONTACTS


def _slide_closing(prs, params: dict, credits: list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Фото фон или тёмная заливка
    img, attr = fetch_photo('Japan zen minimalist garden', 'Japan morning mist landscape')
    if img:
        _photo_bg(slide, img, C.DARK, overlay_opacity=0.68)
        if attr:
            credits.append(attr)
    else:
        _bg(slide, C.DARK)

    # Декоративные линии
    _rect(slide, M, 0.55, W - M * 2, 0.03, C.GOLD)
    _rect(slide, M, H - 0.65, W - M * 2, 0.03, C.GOLD)

    # Главный заголовок
    _txt(slide, 'Создадим это\nвместе?',
         M, 1.1, 8.0, 2.4,
         size=46, bold=True, color=C.WHITE, font='Georgia')

    # Подзаголовок
    _txt(slide, f'Готовы разработать программу для {params["company_name"]}',
         M, 3.6, 9.0, 0.65,
         size=16, italic=True, color=C.GOLD, font='Georgia')

    # Контакты
    c = _get_contacts()
    contact_lines = []
    if c.get('name'):
        contact_lines.append(c['name'])
    if c.get('title'):
        contact_lines.append(c['title'])
    if c.get('email'):
        contact_lines.append(c['email'])
    if c.get('phone'):
        contact_lines.append(c['phone'])
    if c.get('web'):
        contact_lines.append(c['web'])

    _multiline(slide, contact_lines,
               M, 4.45, 7.0, 1.8,
               size=12, color=RGBColor(0xCC, 0xC9, 0xC2),
               font='Arial', line_spacing=1.7)

    # Дисклеймер
    _txt(slide,
         'Международные перелёты не включены. Цены предварительные, подтверждаются при бронировании.',
         M, H - 1.05, 10.0, 0.45,
         size=7.5, italic=True, color=C.STONE)

    _logo(slide, right_align=True)


# ── Слайд: Фотокредиты (скрытый) ─────────────────────────────────────────────

def _slide_credits(prs, credits: list):
    if not credits:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C.WHITE)
    _label(slide, 'Photo Credits', M, M, 5.0, C.STONE)
    unique = list(dict.fromkeys(credits))
    _multiline(slide, unique,
               M, M + 0.4, W - M * 2, H - M * 2,
               size=9, color=C.STONE, font='Arial', line_spacing=1.6)
    # Скрываем слайд (атрибут show="0")
    slide._element.set('show', '0')


# ── Главная функция ───────────────────────────────────────────────────────────

def create_ppt(params: dict, content: dict, services: list) -> bytes:
    reset_session()  # каждая новая презентация — свежий пул фото

    season, _, season_photo = detect_season(params.get('dates', ''))

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    credits: list = []

    _slide_title(prs, params, credits, season_photo=season_photo)
    _slide_why_japan(prs, params, season, credits)
    _slide_concept(prs, content, credits)
    _slide_overview(prs, params, content)

    for day_data in content.get('days', []):
        _slide_day(prs, day_data, credits, season)

    for i, hl in enumerate(content.get('highlights', [])):
        _slide_highlight(prs, hl, credits, index=i)

    _slide_hotels(prs, params, content, credits)
    _slide_closing(prs, params, credits)
    _slide_credits(prs, credits)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
