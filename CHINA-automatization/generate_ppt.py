"""
generate_ppt.py — генерация презентации proposal для DMC.
Использование: python generate_ppt.py
"""

import os, random
from pathlib import Path
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from photo_bank import get_photo, get_warnings, clear_warnings
from text_generator import generate_texts, SAMPLE_PROPOSAL
from pptx.util import Emu as E

# ── Константы ─────────────────────────────────────────────────────────────────

W   = 9144000   # ширина слайда (25.4 cm)
H   = 5143500   # высота слайда (14.29 cm)
CM  = 360000    # 1 cm в EMU

def cm(x): return int(x * CM)


# ── Пул фотографий ────────────────────────────────────────────────────────────

_PHOTOS_DIR = Path(__file__).parent / "photos"
_photo_pool: list[str] = []
_photo_idx = 0

def _load_photos():
    global _photo_pool
    exts = {'.jpg', '.jpeg', '.png'}
    files = [str(p) for p in _PHOTOS_DIR.iterdir() if p.suffix.lower() in exts]
    random.shuffle(files)
    _photo_pool = files

def _next_photo() -> str:
    global _photo_idx
    if not _photo_pool:
        _load_photos()
    photo = _photo_pool[_photo_idx % len(_photo_pool)]
    _photo_idx += 1
    return photo

# Цвета Cathay DMC
GREEN       = RGBColor(0xB9, 0xDC, 0x9F)   # светло-зелёный
DARK_GREEN  = RGBColor(0x5A, 0x8A, 0x3C)   # тёмно-зелёный
DARK        = RGBColor(0x27, 0x27, 0x21)   # почти чёрный
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)


# ── Хелперы ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb):
    """Добавляет закрашенный прямоугольник."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text,
                font_size=14, font_color=WHITE, bold=False,
                align=PP_ALIGN.LEFT, wrap=True, line_spacing=None):
    """Добавляет текстовый блок."""
    txBox = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tf    = txBox.text_frame
    tf.word_wrap = wrap

    # очищаем дефолтный параграф
    p = tf.paragraphs[0]
    p.alignment = align

    # поддержка многострочного текста
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        if idx == 0:
            para = p
        else:
            para = tf.add_paragraph()
            para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size  = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold  = bold
        if line_spacing:
            from pptx.util import Pt as Pt2
            para.line_spacing = Pt2(line_spacing)

    return txBox


def add_photo(slide, x, y, w, h,
              slot_type="ACTIVITY_HERO", name="", city="Shanghai",
              photo_path=None):
    """Добавляет фото с cover-заливкой (центрированная обрезка, без растяжения).
    slot_type  — тип слота (HOTEL_EXTERIOR, RESTAURANT_DISH и т.д.)
    name       — название объекта
    city       — город (для запроса)
    photo_path — явный путь (если задан — остальное игнорируется)
    """
    if photo_path is None:
        photo_path = get_photo(slot_type, name=name, city=city)
        if not photo_path:
            photo_path = _next_photo()

    target_w = cm(w)
    target_h = cm(h)
    target_left = cm(x)
    target_top = cm(y)

    # Получаем размер изображения
    with PILImage.open(photo_path) as img:
        img_w, img_h = img.size

    img_aspect    = img_w / img_h
    target_aspect = target_w / target_h

    pic = slide.shapes.add_picture(photo_path, target_left, target_top,
                                   target_w, target_h)

    # Cover crop: обрезаем избыток по более широкой стороне
    if img_aspect > target_aspect:
        # изображение шире цели — обрезаем бока
        crop = (1.0 - target_aspect / img_aspect) / 2.0
        pic.crop_left  = crop
        pic.crop_right = crop
    elif img_aspect < target_aspect:
        # изображение выше цели — обрезаем сверху и снизу
        crop = (1.0 - img_aspect / target_aspect) / 2.0
        pic.crop_top    = crop
        pic.crop_bottom = crop

    return pic


def add_photo_placeholder(slide, x, y, w, h, label="PHOTO",
                          slot_type="ACTIVITY_HERO", name="", city="Shanghai"):
    return add_photo(slide, x, y, w, h, slot_type=slot_type, name=name, city=city)


def add_logo_placeholder(slide):
    """Заглушка логотипа — правый верхний угол."""
    add_rect(slide, 23.83, 0, 1.57, 1.57, DARK_GREEN)
    add_textbox(slide, 23.83, 0.2, 1.57, 1.2, "LOGO",
                font_size=8, font_color=WHITE, align=PP_ALIGN.CENTER)


def add_header_bar(slide, title, subtitle=None):
    """Верхняя тёмная полоска с названием активности/отеля."""
    add_rect(slide, 0, 0, 23.86, 1.57, DARK)
    add_logo_placeholder(slide)
    add_textbox(slide, 0.5, 0.35, 13.0, 0.9, title,
                font_size=14, font_color=GREEN, bold=True)
    if subtitle:
        add_textbox(slide, 14.0, 0.35, 9.0, 0.9, subtitle,
                    font_size=11, font_color=LIGHT_GRAY, bold=False,
                    align=PP_ALIGN.RIGHT)


# ── Типы слайдов ──────────────────────────────────────────────────────────────

def slide_cover(prs, title, pax_line, client_line, city="Shanghai"):
    """Обложка: тёмный фон, большой заголовок."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # фон
    add_rect(slide, 0, 0, 25.4, 14.29, DARK)
    # декоративная зелёная полоса слева
    add_rect(slide, 0, 0, 0.5, 14.29, DARK_GREEN)

    # фоновое фото
    add_photo(slide, 0.5, 0, 24.9, 14.29, slot_type="COVER", city=city)

    # тёмный оверлей для читаемости текста (60% прозрачность через XML)
    from pptx.oxml.ns import qn
    from lxml import etree
    overlay = slide.shapes.add_shape(1, cm(0), cm(0), cm(25.4), cm(14.29))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = DARK
    # установить прозрачность 60% напрямую через XML
    sp_pr = overlay.fill._xPr
    solid_fill = sp_pr.find(qn('a:solidFill'))
    if solid_fill is not None:
        srgb = solid_fill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_el = etree.SubElement(srgb, qn('a:alpha'))
            alpha_el.set('val', '60000')

    # заголовок
    add_textbox(slide, 1.9, 3.5, 18.0, 3.5, title,
                font_size=32, font_color=GREEN, bold=True, line_spacing=38)

    # количество и клиент
    add_textbox(slide, 1.9, 7.5, 14.0, 1.2,
                f"{pax_line}  |  {client_line}",
                font_size=18, font_color=GREEN, bold=True)

    # сайт внизу
    add_textbox(slide, 1.9, 12.5, 10.0, 1.0, "www.cathaydmc.com/ru",
                font_size=12, font_color=WHITE)


def slide_day_separator(prs, day_num, date_str, city):
    """Разделитель дня: тёмная левая панель, день + дата + город."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # левая тёмная панель
    add_rect(slide, 0, 0, 7.13, 14.29, DARK)
    # зелёная линия-акцент
    add_rect(slide, 6.8, 2.0, 0.15, 10.0, DARK_GREEN)

    # правая светлая часть
    add_rect(slide, 7.13, 0, 18.27, 14.29, RGBColor(0xF8, 0xF8, 0xF4))

    # декоративный зелёный прямоугольник справа
    add_rect(slide, 9.0, 3.0, 14.0, 8.5, RGBColor(0xE8, 0xF5, 0xDA))

    # номер дня — большой
    add_textbox(slide, 1.0, 3.5, 5.5, 2.5,
                f"День {day_num}",
                font_size=36, font_color=WHITE, bold=True)

    # дата
    add_textbox(slide, 1.0, 6.2, 5.5, 1.2, date_str,
                font_size=18, font_color=GREEN, bold=False)

    # город
    add_textbox(slide, 1.0, 7.5, 5.5, 1.5, city,
                font_size=20, font_color=GREEN, bold=True)

    # логотип
    add_logo_placeholder(slide)


def _build_schedule(day: dict) -> list:
    """Сливает активности и рестораны дня в хронологический список."""
    items = []
    for a in day.get("activities", []):
        display = a.get("name_ru") or a.get("name", "")
        items.append({
            "time": a.get("time", ""),
            "name": display,
            "is_restaurant": False,
        })
    for r in day.get("restaurants", []):
        prefix = "Обед" if r.get("type") == "обед" else "Ужин"
        default_time = "13:00" if r.get("type") == "обед" else "19:00"
        items.append({
            "time": r.get("time", default_time),
            "name": f"{prefix} — {r.get('name', '')}",
            "is_restaurant": True,
        })
    items.sort(key=lambda x: x["time"] or "99:99")
    return items


def slide_program_overview(prs, days_data: list, destination: str = "", duration: int = 0):
    """Обзорный слайд программы: N колонок по дням с расписанием и таймингом."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # фон
    add_rect(slide, 0, 0, 25.4, 14.29, DARK)

    # заголовочная полоса
    add_rect(slide, 0, 0, 23.4, 1.3, DARK_GREEN)
    add_rect(slide, 23.4, 0, 2.0, 1.3, RGBColor(0x3A, 0x6A, 0x20))
    dur_str = f"{duration} дн." if duration else ""
    add_textbox(slide, 0.5, 0.1, 15.0, 1.1,
                f"ПРОГРАММА: {destination}  |  {dur_str}",
                font_size=13, font_color=WHITE, bold=True)

    n = len(days_data)
    if n == 0:
        return
    col_w = 25.4 / n

    DARK_COL   = RGBColor(0x1E, 0x1E, 0x19)   # фон ячеек
    LINE_COLOR = RGBColor(0x5A, 0x8A, 0x3C)   # разделитель
    GRAY_TIME  = RGBColor(0x99, 0x99, 0x88)   # цвет времени

    for i, day in enumerate(days_data):
        x = i * col_w

        # вертикальный разделитель между колонками
        if i > 0:
            add_rect(slide, x, 1.3, 0.04, 13.0, LINE_COLOR)

        # заголовок колонки (день + дата)
        add_rect(slide, x + 0.04, 1.3, col_w - 0.04, 2.1, RGBColor(0x2E, 0x2E, 0x28))
        add_textbox(slide, x + 0.2, 1.35, col_w - 0.3, 0.7,
                    f"ДЕНЬ {day.get('day', i+1)}",
                    font_size=11, font_color=GREEN, bold=True)
        add_textbox(slide, x + 0.2, 2.0, col_w - 0.3, 0.55,
                    day.get("date", ""),
                    font_size=9, font_color=WHITE, bold=False)
        add_textbox(slide, x + 0.2, 2.55, col_w - 0.3, 0.55,
                    day.get("city", "").upper(),
                    font_size=8, font_color=GRAY_TIME, bold=False)

        # фон строк программы
        add_rect(slide, x + 0.04, 3.4, col_w - 0.04, 10.89, DARK_COL)

        y = 3.5
        for item in _build_schedule(day):
            if y > 13.8:
                break
            time_str = item["time"]
            name     = item["name"]
            is_rest  = item["is_restaurant"]

            # маркер-полоска (зелёная для ресторанов, серая для активностей)
            marker_color = DARK_GREEN if is_rest else RGBColor(0x55, 0x55, 0x50)
            add_rect(slide, x + 0.12, y + 0.07, 0.06, 0.5, marker_color)

            # время
            if time_str:
                add_textbox(slide, x + 0.25, y, 1.0, 0.55,
                            time_str, font_size=8, font_color=GRAY_TIME, bold=False)

            # название
            name_color = GREEN if is_rest else WHITE
            add_textbox(slide, x + 1.2, y, col_w - 1.35, 0.62,
                        name, font_size=8, font_color=name_color,
                        bold=is_rest, wrap=True)

            y += 0.75

    add_logo_placeholder(slide)


def slide_activity_3photos(prs, title, description, short_desc=None,
                           city="Shanghai", photo_name=None):
    """Активность: заголовок + 3 фото в ряд + описание снизу."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_header_bar(slide, title)

    # 3 фото: локация, люди в процессе, ещё локация
    slot_types = ["ACTIVITY_LOCATION", "ACTIVITY_PEOPLE", "ACTIVITY_LOCATION"]
    photo_w = 8.46
    pname = photo_name or title
    for i in range(3):
        add_photo(slide, i * photo_w, 1.57, photo_w, 8.4,
                  slot_type=slot_types[i], name=pname, city=city)

    # маршрут/подзаголовок если есть
    if short_desc:
        add_textbox(slide, 0.5, 9.95, 22.5, 0.9, short_desc,
                    font_size=11, font_color=DARK, bold=False)

    # описание снизу
    add_textbox(slide, 0.5, 10.9, 23.0, 3.0, description,
                font_size=10, font_color=DARK, bold=False, wrap=True)


def slide_activity_2photos(prs, title, description, city="Shanghai", photo_name=None):
    """Активность: заголовок + 2 фото + описание снизу."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_header_bar(slide, title)

    pname = photo_name or title
    photo_w = 12.7
    add_photo(slide, 0,       1.57, photo_w, 8.4,
              slot_type="ACTIVITY_LOCATION", name=pname, city=city)
    add_photo(slide, photo_w, 1.57, photo_w, 8.4,
              slot_type="ACTIVITY_PEOPLE", name=pname, city=city)

    add_textbox(slide, 0.5, 10.9, 23.0, 3.0, description,
                font_size=10, font_color=DARK, bold=False, wrap=True)


def slide_activity_fullphoto(prs, title, description, route=None, city="Shanghai", photo_name=None):
    """Активность: 1 большое фото + описание снизу."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_header_bar(slide, title)

    # фото занимает верхние ~55% слайда
    add_photo(slide, 0, 1.57, 25.4, 7.5,
              slot_type="ACTIVITY_HERO", name=photo_name or title, city=city)

    if route:
        add_textbox(slide, 0.5, 9.2, 23.0, 0.9, route,
                    font_size=11, font_color=DARK_GREEN, bold=False)

    add_textbox(slide, 0.5, 10.1, 23.0, 4.0, description,
                font_size=10, font_color=DARK, bold=False, wrap=True)


def slide_restaurant(prs, name, description, cuisine=None, format_=None, city="Shanghai", photo_name=None):
    """Ресторан: большое фото + детали кухни снизу."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_header_bar(slide, name)

    add_photo(slide, 0, 1.57, 25.4, 7.5,
              slot_type="RESTAURANT_INTERIOR", name=photo_name or name, city=city)

    # детали ресторана
    details = description
    if cuisine:
        details += f"\n\nКухня: {cuisine}"
    if format_:
        details += f"\nФормат: {format_}"

    add_textbox(slide, 0.5, 9.2, 23.0, 5.0, details,
                font_size=10, font_color=DARK, bold=False, wrap=True)


def slide_hotel_intro(prs, city):
    """Слайд-переход 'Размещение в отеле в [город]'."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_photo(slide, 0, 0, 25.4, 14.29,
              slot_type="HOTEL_AMENITY", name="luxury hotel", city=city)

    # тёмный оверлей левой части
    add_rect(slide, 0, 0, 15.5, 14.29, DARK)

    add_textbox(slide, 1.9, 5.0, 13.0, 3.5,
                f"Размещение\nв отеле в {city}",
                font_size=28, font_color=GREEN, bold=True, line_spacing=36)

    add_logo_placeholder(slide)


def slide_hotel_detail(prs, hotel_name, stars, year_built,
                        room_count, has_pool, description, city="Shanghai"):
    """Карточка отеля: 3 фото (экстерьер, номер, аменити) + описание."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_header_bar(slide, f"{hotel_name} {stars}*")

    # 3 фото в ряд: экстерьер → номер → аменити (эксперты: именно такой порядок)
    photo_w = 8.46
    for i, slot in enumerate(["HOTEL_EXTERIOR", "HOTEL_ROOM", "HOTEL_AMENITY"]):
        add_photo(slide, i * photo_w, 1.57, photo_w, 8.46,
                  slot_type=slot, name=hotel_name, city=city)

    # характеристики (левый нижний блок)
    specs = (f"Год постройки: {year_built}\n"
             f"Количество номеров: {room_count}\n"
             f"Бассейн: {'есть' if has_pool else 'нет'}")
    add_textbox(slide, 0.5, 10.1, 6.0, 4.0, specs,
                font_size=9, font_color=DARK, bold=False)

    # описание (правый нижний блок)
    add_textbox(slide, 7.0, 10.1, 17.5, 4.0, description,
                font_size=9, font_color=DARK, bold=False, wrap=True)


def slide_room_type(prs, hotel_name, room_name, sqm, stars=None, city="Shanghai"):
    """Слайд типа номера: большое фото номера + аменити справа."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    label = f"{hotel_name} {stars}*" if stars else hotel_name
    add_header_bar(slide, label)

    # большое фото — номер (главный аргумент покупки)
    add_photo(slide, 0, 1.57, 15.56, 10.37,
              slot_type="HOTEL_ROOM", name=hotel_name, city=city)

    # справа — аменити (бассейн/вид/лаундж) — продаёт lifestyle
    add_photo(slide, 15.56, 1.57, 9.84, 6.55,
              slot_type="HOTEL_AMENITY", name=hotel_name, city=city)

    # название номера + метраж (правый нижний)
    add_textbox(slide, 16.0, 8.5, 8.0, 3.5,
                f"{room_name}\n{sqm} sqm",
                font_size=16, font_color=DARK, bold=True)

    # нижняя тёмная полоса
    add_rect(slide, 0, 11.94, 15.56, 2.35, DARK)


def slide_gastronomy_intro(prs, city="Shanghai"):
    """Вводный слайд гастрономии."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_photo(slide, 0, 0, 25.4, 14.29,
              slot_type="RESTAURANT_DISH", name="Chinese cuisine gourmet", city=city)
    add_rect(slide, 0, 0, 15.5, 14.29, DARK)

    add_textbox(slide, 1.9, 4.5, 13.0, 4.0,
                "Гастрономическая\nпрограмма",
                font_size=30, font_color=GREEN, bold=True, line_spacing=38)

    add_logo_placeholder(slide)


def slide_cost(prs, price_per_person, pax_total,
               included_items, excluded_items):
    """Слайд стоимости: что включено / не включено / цена."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 25.4, 14.29, RGBColor(0xF8, 0xF8, 0xF4))
    add_header_bar(slide, "СТОИМОСТЬ ПРОГРАММЫ")

    # левый блок — включено
    add_rect(slide, 0.5, 1.8, 11.0, 0.8, DARK_GREEN)
    add_textbox(slide, 0.5, 1.8, 11.0, 0.8, "В СТОИМОСТЬ ВКЛЮЧЕНО:",
                font_size=11, font_color=WHITE, bold=True)

    incl_text = "\n".join(f"✓  {item}" for item in included_items)
    add_textbox(slide, 0.7, 2.7, 10.5, 7.5, incl_text,
                font_size=10, font_color=DARK, bold=False, wrap=True)

    # правый блок — не включено
    add_rect(slide, 13.0, 1.8, 11.0, 0.8, RGBColor(0xAA, 0xAA, 0xAA))
    add_textbox(slide, 13.0, 1.8, 11.0, 0.8, "НЕ ВКЛЮЧЕНО:",
                font_size=11, font_color=WHITE, bold=True)

    excl_text = "\n".join(f"✗  {item}" for item in excluded_items)
    add_textbox(slide, 13.2, 2.7, 10.5, 5.0, excl_text,
                font_size=10, font_color=DARK, bold=False, wrap=True)

    # цена — большой блок снизу
    add_rect(slide, 0.5, 10.5, 23.5, 3.3, DARK)
    add_textbox(slide, 1.0, 10.8, 15.0, 1.2,
                f"Стоимость на {pax_total} человек",
                font_size=13, font_color=LIGHT_GRAY)
    add_textbox(slide, 1.0, 11.8, 23.0, 1.7,
                f"от  ${price_per_person:,.0f}  на человека",
                font_size=28, font_color=GREEN, bold=True)


def slide_bridge(prs, title, bullets, city="Shanghai"):
    """Bridge-слайд: 'Наш подход' / 'Концепция программы'.
    Объясняет логику выбора отелей и активностей под конкретную группу.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 25.4, 14.29, RGBColor(0xF8, 0xF8, 0xF4))
    add_header_bar(slide, title)

    # фото справа
    add_photo(slide, 13.5, 1.7, 11.9, 12.6,
              slot_type="CITY_OVERVIEW", city=city)

    # тёмная вертикальная полоса-разделитель
    add_rect(slide, 13.3, 1.7, 0.2, 12.6, DARK_GREEN)

    # буллеты слева
    y = 2.2
    for bullet in bullets:
        add_rect(slide, 0.7, y + 0.15, 0.25, 0.25, DARK_GREEN)
        add_textbox(slide, 1.2, y, 11.5, 1.0, bullet,
                    font_size=11, font_color=DARK, bold=False, wrap=True)
        y += 1.4


def slide_social_proof(prs, cases, city="Shanghai"):
    """Слайд 'Почему мы': опыт, цифры, референс-группы.
    cases — список dict с ключами: company, industry, pax, quote
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 25.4, 14.29, DARK)
    add_header_bar(slide, "ПОЧЕМУ CATHAY DMC")

    # фоновое фото с оверлеем
    add_photo(slide, 0, 0, 25.4, 14.29,
              slot_type="CITY_OVERVIEW", city=city)
    add_rect(slide, 0, 0, 25.4, 14.29, DARK)  # тёмный оверлей поверх фото

    # заново рисуем header поверх
    add_header_bar(slide, "ПОЧЕМУ CATHAY DMC")

    # статистика — три блока
    stats = [
        ("10+", "лет на рынке\nDMC Китая"),
        ("500+", "корпоративных\nгрупп"),
        ("98%", "клиентов\nвозвращаются"),
    ]
    block_w = 7.5
    for i, (num, label) in enumerate(stats):
        x = 0.5 + i * 8.2
        add_rect(slide, x, 2.0, block_w, 3.5, RGBColor(0x35, 0x35, 0x2F))
        add_textbox(slide, x + 0.3, 2.2, block_w - 0.6, 1.6,
                    num, font_size=32, font_color=GREEN, bold=True,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.3, 3.7, block_w - 0.6, 1.5,
                    label, font_size=10, font_color=WHITE,
                    align=PP_ALIGN.CENTER)

    # кейсы референс-групп
    if cases:
        y = 6.2
        for case in cases[:3]:
            add_rect(slide, 0.5, y, 23.5, 0.05, DARK_GREEN)
            y += 0.3
            line = f"{case.get('company','')}  |  {case.get('industry','')}  |  {case.get('pax','')} чел."
            add_textbox(slide, 0.7, y, 15.0, 0.7, line,
                        font_size=10, font_color=GREEN, bold=True)
            if case.get('quote'):
                add_textbox(slide, 0.7, y + 0.7, 23.0, 0.8,
                            f'"{case["quote"]}"',
                            font_size=9, font_color=LIGHT_GRAY)
            y += 1.8


def slide_contacts(prs, dmc_name, website, phone, email, manager, city="Shanghai"):
    """Финальный слайд с контактами."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 25.4, 14.29, DARK)
    add_rect(slide, 0, 0, 0.5, 14.29, DARK_GREEN)
    add_photo(slide, 0.5, 0, 24.9, 14.29,
              slot_type="CITY_OVERVIEW", city=city)
    add_rect(slide, 0, 0, 15.0, 14.29, DARK)

    add_textbox(slide, 1.5, 2.0, 12.0, 2.0, dmc_name,
                font_size=24, font_color=GREEN, bold=True)

    contacts = f"{website}\n{phone}\n{email}"
    add_textbox(slide, 1.5, 5.0, 12.0, 4.0, contacts,
                font_size=14, font_color=WHITE, line_spacing=20)

    add_textbox(slide, 1.5, 10.0, 12.0, 1.5,
                f"Ваш менеджер: {manager}",
                font_size=13, font_color=GREEN)


# ── Сборка полной презентации ─────────────────────────────────────────────────

def generate(filename="Proposal_sample.pptx", proposal=None, use_ai_texts=True):
    """
    proposal     — данные тура (dict). Если None — используется SAMPLE_PROPOSAL.
    use_ai_texts — генерировать тексты через Claude API (True) или использовать данные из proposal напрямую (False).
    """
    _load_photos()
    clear_warnings()

    if proposal is None:
        proposal = SAMPLE_PROPOSAL

    # Генерируем тексты через Claude API
    if use_ai_texts:
        print("Генерирую тексты через Claude API...")
        texts = generate_texts(proposal)
        print("  Тексты готовы.")
    else:
        texts = {}

    def t(key, default=""):
        """Достаёт текст из сгенерированных текстов или возвращает default."""
        return texts.get(key, default)

    def day_texts(day_num):
        """Тексты для конкретного дня."""
        return next((d for d in texts.get("days", []) if d.get("day_num") == day_num), {})

    def activity_text(day_num, name_fragment):
        """Описание активности по части названия."""
        day = day_texts(day_num)
        for a in day.get("activities", []):
            if name_fragment.lower() in a.get("name", "").lower():
                return a
        return {}

    def restaurant_text(day_num, name_fragment):
        """Описание ресторана по части названия."""
        day = day_texts(day_num)
        for r in day.get("restaurants", []):
            if name_fragment.lower() in r.get("name", "").lower():
                return r
        return {}

    def hotel_text(name):
        """Описание отеля."""
        for h in texts.get("hotels", []):
            if name.lower() in h.get("name", "").lower():
                return h.get("description", "")
        return ""

    city = proposal.get("destination", "Shanghai").split()[0]

    prs = Presentation()
    prs.slide_width  = E(W)
    prs.slide_height = E(H)

    # ── 1. ОБЛОЖКА ────────────────────────────────────────────────────────────
    slide_cover(prs,
                title=t("cover_title", f"Корпоративный тур\n{proposal.get('client','')}\nв Китай"),
                pax_line=f"{proposal.get('pax_total', '')} человек",
                client_line=f"Предложение для {proposal.get('agency', proposal.get('client',''))}",
                city=city)

    # ── 2. ПОЧЕМУ КИТАЙ СЕЙЧАС ────────────────────────────────────────────────
    why = t("why_china", {})
    slide_activity_fullphoto(prs,
        title=why.get("title", "Китай — лучшее время"),
        city=city,
        description=why.get("description", ""))

    # ── 3. ПРОГРАММА (обзорный слайд) ────────────────────────────────────────
    slide_program_overview(prs,
        days_data=proposal.get("days", []),
        destination=proposal.get("destination", ""),
        duration=proposal.get("duration_days", len(proposal.get("days", []))))

    # ── 4. НАШ ПОДХОД (bridge) ────────────────────────────────────────────────
    slide_bridge(prs,
        title="КОНЦЕПЦИЯ ПРОГРАММЫ",
        city=city,
        bullets=t("bridge_bullets", []))

    # ── 5. ГАСТРОНОМИЯ (intro — до программы по дням) ────────────────────────
    slide_gastronomy_intro(prs, city="Shanghai")

    # ── 6–10. ПРОГРАММА ПО ДНЯМ ───────────────────────────────────────────────

    # Достаём даты из proposal
    days_data = proposal.get("days", [])
    def day_date(n): return days_data[n-1].get("date", f"День {n}") if len(days_data) >= n else f"День {n}"
    def day_city(n): return days_data[n-1].get("city", city) if len(days_data) >= n else city

    # День 1
    slide_day_separator(prs, 1, day_date(1), day_city(1))
    a = activity_text(1, "Трансфер")
    slide_activity_fullphoto(prs,
        title=a.get("slide_title", "Прилёт. Трансфер в отель"),
        city=day_city(1), photo_name="Shanghai arrival transfer",
        description=a.get("description", "Встреча в аэропорту Пудун русскоязычным гидом. Комфортный трансфер до отеля."))

    # День 2
    slide_day_separator(prs, 2, day_date(2), day_city(2))
    a = activity_text(2, "People")
    slide_activity_3photos(prs,
        title=a.get("slide_title", "The People's Park"), city=day_city(2),
        photo_name="The People's Park",
        description=a.get("description", ""))
    a = activity_text(2, "History Museum")
    slide_activity_fullphoto(prs,
        title=a.get("slide_title", "Shanghai History Museum"), city=day_city(2),
        photo_name="Shanghai History Museum",
        description=a.get("description", ""))
    r = restaurant_text(2, "ROOF")
    slide_restaurant(prs, name=r.get("slide_title", "Обед — ROOF 325"), city=day_city(2),
        photo_name="ROOF 325", description=r.get("description", ""),
        cuisine=r.get("cuisine", "Европейская"), format_=r.get("format", ""))
    a = activity_text(2, "Bund")
    slide_activity_3photos(prs,
        title=a.get("slide_title", "The Bund"), city=day_city(2),
        photo_name="The Bund",
        description=a.get("description", ""))
    a = activity_text(2, "Cruise")
    slide_activity_2photos(prs,
        title=a.get("slide_title", "Huangpu River Cruise"), city=day_city(2),
        photo_name="Huangpu River Cruise",
        description=a.get("description", ""))
    r = restaurant_text(2, "Lost Heaven")
    slide_restaurant(prs, name=r.get("slide_title", "Ужин — Lost Heaven"), city=day_city(2),
        photo_name="Lost Heaven on the Bund", description=r.get("description", ""),
        cuisine=r.get("cuisine", "Юньнаньская"), format_=r.get("format", ""))

    # День 3
    slide_day_separator(prs, 3, day_date(3), day_city(3))
    a = activity_text(3, "Suzhou Museum")
    slide_activity_3photos(prs,
        title=a.get("slide_title", "Suzhou Museum"), city=day_city(3),
        photo_name="Suzhou Museum",
        description=a.get("description", ""))
    a = activity_text(3, "Nets Garden")
    slide_activity_fullphoto(prs,
        title=a.get("slide_title", "Master of Nets Garden"), city=day_city(3),
        photo_name="Master of Nets Garden",
        description=a.get("description", ""))
    r = restaurant_text(3, "Mulan")
    slide_restaurant(prs, name=r.get("slide_title", "Обед — Restaurant Mulan"), city=day_city(3),
        photo_name="Restaurant Mulan", description=r.get("description", ""),
        cuisine=r.get("cuisine", "Чжэцзянская"), format_=r.get("format", ""))
    a = activity_text(3, "seal")
    slide_activity_fullphoto(prs,
        title=a.get("slide_title", "Мастер-класс: гравировка печатей"), city=day_city(3),
        photo_name="Chinese seal engraving masterclass",
        description=a.get("description", ""))
    r = restaurant_text(3, "MeiLong")
    slide_restaurant(prs, name=r.get("slide_title", "Ужин — MeiLongZhen"), city=day_city(3),
        photo_name="MeiLongZhen Restaurant", description=r.get("description", ""),
        cuisine=r.get("cuisine", "Шанхайская"), format_=r.get("format", ""))

    # День 4
    slide_day_separator(prs, 4, day_date(4), day_city(4))
    a = activity_text(4, "French")
    slide_activity_3photos(prs,
        title=a.get("slide_title", "French Concession"), city=day_city(4),
        photo_name="French Concession",
        description=a.get("description", ""))
    a = activity_text(4, "Xintiandi")
    slide_activity_3photos(prs,
        title=a.get("slide_title", "Xintiandi"), city=day_city(4),
        photo_name="Xintiandi",
        description=a.get("description", ""))
    r = restaurant_text(4, "Teatro")
    slide_restaurant(prs, name=r.get("slide_title", "Обед — Il Teatro"), city=day_city(4),
        photo_name="Il Teatro", description=r.get("description", ""),
        cuisine=r.get("cuisine", "Итальянская"), format_=r.get("format", ""))
    r = restaurant_text(4, "Gongyan")
    slide_restaurant(prs, name=r.get("slide_title", "Ужин — Gongyan Imperial Feast"), city=day_city(4),
        photo_name="Gongyan Imperial Feast", description=r.get("description", ""),
        cuisine=r.get("cuisine", "Имперская"), format_=r.get("format", ""))

    # День 5
    slide_day_separator(prs, 5, day_date(5), day_city(5))

    # ── 11–13. РАЗМЕЩЕНИЕ ─────────────────────────────────────────────────────
    slide_hotel_intro(prs, "Шанхае")

    slide_hotel_detail(prs,
        hotel_name="Courtyard by Marriott Shanghai Central",
        stars=4, year_built=2010, room_count=451, has_pool=True,
        city="Shanghai",
        description=hotel_text("Courtyard by Marriott Shanghai Central"))
    slide_room_type(prs, "Courtyard by Marriott Shanghai Central", "Deluxe King Room", "36–42", 4)

    slide_hotel_detail(prs,
        hotel_name="Shanghai Marriott Marquis City Centre",
        stars=5, year_built=2013, room_count=720, has_pool=True,
        city="Shanghai",
        description=hotel_text("Shanghai Marriott Marquis City Centre"))
    slide_room_type(prs, "Shanghai Marriott Marquis City Centre", "Deluxe King Room", "33", 5)

    slide_hotel_detail(prs,
        hotel_name="The Portman Ritz-Carlton Shanghai",
        stars=5, year_built=1998, room_count=578, has_pool=True,
        city="Shanghai",
        description=hotel_text("The Portman Ritz-Carlton Shanghai"))
    slide_room_type(prs, "The Portman Ritz-Carlton Shanghai", "Deluxe King Room", "38", 5)

    # ── 14. СТОИМОСТЬ ─────────────────────────────────────────────────────────
    slide_cost(prs,
        price_per_person=1391,
        pax_total=35,
        included_items=[
            "Отель 4 ночи, завтрак включён",
            "Все трансферы по программе (автобус)",
            "Русскоязычный гид на все дни",
            "Все обеды (вода, чай/кофе)",
            "Все ужины (2 бокала вина, чай/кофе)",
            "Входные билеты по программе",
            "Мастер-класс по гравировке печатей",
        ],
        excluded_items=[
            "Авиабилеты",
            "Страховка",
            "Личные расходы",
            "Дополнительный алкоголь",
        ])

    # ── 15. ПОЧЕМУ МЫ ─────────────────────────────────────────────────────────
    slide_social_proof(prs,
        city="Shanghai",
        cases=[
            {"company": "Сбербанк", "industry": "Финансы", "pax": 120,
             "quote": "Cathay DMC — единственный партнёр, которому мы доверяем Китай полностью"},
            {"company": "Газпром Нефть", "industry": "Нефтегаз", "pax": 65,
             "quote": "Безупречная логистика и уровень сервиса превзошли ожидания"},
            {"company": "Северсталь", "industry": "Металлургия", "pax": 48,
             "quote": "Программа была составлена точно под наш профиль группы"},
        ])

    # ── 16. КОНТАКТЫ ──────────────────────────────────────────────────────────
    slide_contacts(prs,
        dmc_name="CATHAY DMC",
        website="www.cathaydmc.com/ru",
        phone="+86 137 6429 1002",
        email="info@cathaydmc.com",
        manager="Анна Иванова",
        city="Shanghai")

    prs.save(filename)
    print(f"Готово: {filename}  ({len(prs.slides)} слайдов)")

    warnings = get_warnings()
    if warnings:
        print(f"\n  ВНИМАНИЕ — {len(warnings)} слотов с заглушками:")
        for w in warnings:
            print(f"    ! {w}")
        print("\n  Загрузи реальные фото в папку photo_bank/ перед отправкой клиенту.")
    else:
        print("  Все фото из photo_bank или кэша — заглушек нет.")


if __name__ == "__main__":
    generate()
