# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image

BASE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(BASE, 'assets')

INK = RGBColor(0x22, 0x1A, 0x12)
PARCHMENT = RGBColor(0xFB, 0xF3, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

ACCENTS = {
    'rust':    {'light': RGBColor(0xE8, 0x5A, 0x2C), 'dark': RGBColor(0xB7, 0x39, 0x1A), 'darker': RGBColor(0x8F, 0x2A, 0x12)},
    'mustard': {'light': RGBColor(0xF0, 0xB9, 0x28), 'dark': RGBColor(0xC6, 0x89, 0x0E)},
    'teal':    {'light': RGBColor(0x17, 0xB2, 0x94), 'dark': RGBColor(0x0B, 0x7A, 0x65)},
    'grape':   {'light': RGBColor(0xA9, 0x47, 0x7F), 'dark': RGBColor(0x71, 0x2F, 0x58)},
}

TOP_X = 0.70
BOT_X = 0.48

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def px(pct):
    return Emu(int(SW * pct))


def py(pct):
    return Emu(int(SH * pct))


def set_alpha(shape, pct):
    """pct: 0-100 opacity"""
    srgb = shape.fill.fore_color._xFill
    alpha = etree.SubElement(srgb, qn('a:alpha'))
    alpha.set('val', str(int(pct * 1000)))


def no_line(shape):
    shape.line.fill.background()


def add_bg(slide, accent):
    light = ACCENTS[accent]['light']
    dark = ACCENTS[accent]['dark']
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    base.fill.solid(); base.fill.fore_color.rgb = light
    no_line(base)
    base.shadow.inherit = False
    fb = slide.shapes.build_freeform(px(TOP_X), Emu(0), scale=1.0)
    fb.add_line_segments([
        (SW, Emu(0)),
        (SW, SH),
        (px(BOT_X), SH),
    ], close=True)
    tri = fb.convert_to_shape()
    tri.fill.solid(); tri.fill.fore_color.rgb = dark
    no_line(tri)
    tri.shadow.inherit = False
    return base, tri


def add_card(slide, left, top, width, height, rot=0.0):
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Emu(30000), top + Emu(50000), width, height)
    shadow.adjustments[0] = 0.02
    shadow.fill.solid(); shadow.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    set_alpha(shadow, 30)
    no_line(shadow)
    shadow.shadow.inherit = False
    shadow.rotation = rot

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.02
    card.fill.solid(); card.fill.fore_color.rgb = PARCHMENT
    no_line(card)
    card.shadow.inherit = False
    card.rotation = rot
    return card


def add_text(slide, left, top, width, height, text, font='Manrope', size=14, color=INK,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             rot=0.0, line_spacing=1.0, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.rotation = rot
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing != 1.0:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_price_tag(slide, left, top, accent, price, extra="1 день · от 10 чел"):
    text = f"{price} / чел  ·  {extra}"
    w = Inches(0.16 * len(text) / 1.0 + 1.0)
    h = Inches(0.42)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    pill.adjustments[0] = 0.5
    pill.fill.solid(); pill.fill.fore_color.rgb = ACCENTS[accent]['light']
    no_line(pill)
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(10); tf.margin_right = Pt(10); tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = 'Manrope'; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = PARCHMENT
    return pill, h


def add_beats(slide, left, top, width, beats, accent, size=12.5, gap=Pt(6)):
    tb = slide.shapes.add_textbox(left, top, width, Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for lead, rest in beats:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = gap
        p.line_spacing = 1.05
        rb = p.add_run(); rb.text = '◆  '
        rb.font.size = Pt(size); rb.font.color.rgb = ACCENTS[accent]['light']; rb.font.name = 'Manrope'
        r1 = p.add_run(); r1.text = lead
        r1.font.name = 'Manrope'; r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = INK
        r2 = p.add_run(); r2.text = '  — ' + rest
        r2.font.name = 'Manrope'; r2.font.size = Pt(size); r2.font.bold = False; r2.font.color.rgb = INK
    return tb


def add_postcard(slide, cx, cy, size, img_path, rot):
    frame_w = size
    frame_h = size + Inches(0.3)
    frame_left = cx - frame_w / 2
    frame_top = cy - frame_h / 2

    shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, frame_left + Emu(40000), frame_top + Emu(60000), frame_w, frame_h)
    shadow.fill.solid(); shadow.fill.fore_color.rgb = RGBColor(0, 0, 0)
    set_alpha(shadow, 35)
    no_line(shadow); shadow.shadow.inherit = False
    shadow.rotation = rot

    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, frame_left, frame_top, frame_w, frame_h)
    frame.fill.solid(); frame.fill.fore_color.rgb = WHITE
    no_line(frame); frame.shadow.inherit = False
    frame.rotation = rot

    pad = Inches(0.1)
    photo_size = frame_w - pad * 2
    pic = slide.shapes.add_picture(img_path, frame_left + pad, frame_top + pad, width=photo_size, height=photo_size)
    pic.rotation = rot

    tape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx - Inches(0.32), frame_top - Inches(0.1), Inches(0.64), Inches(0.22))
    tape.fill.solid(); tape.fill.fore_color.rgb = WHITE
    set_alpha(tape, 55)
    no_line(tape); tape.shadow.inherit = False
    tape.rotation = rot - 3


def add_doodle(slide, name, cx, cy, w_in):
    path = os.path.join(ASSETS, name + '.png')
    with Image.open(path) as im:
        iw, ih = im.size
    h_in = w_in * ih / iw
    left = cx - Inches(w_in) / 2
    top = cy - Inches(h_in) / 2
    slide.shapes.add_picture(path, left, top, width=Inches(w_in), height=Inches(h_in))


def add_note(slide, text, left, top, rot):
    add_text(slide, left, top, Inches(4.2), Inches(0.7), text, font='Caveat', size=17,
              color=INK, bold=True, rot=rot, line_spacing=1.0)


# ---------------------------------------------------------------
# SLIDE 1 — COVER
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_bg(s, 'rust')

cover_left = px(0.46)
cover_pic = s.shapes.add_picture(os.path.join(ASSETS, 'ph_cover.jpg'), cover_left, Emu(0),
                                  width=SW - cover_left, height=SH)
tint = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cover_left, Emu(0), SW - cover_left, SH)
tint.fill.solid(); tint.fill.fore_color.rgb = ACCENTS['rust']['darker']
set_alpha(tint, 26)
no_line(tint); tint.shadow.inherit = False

add_doodle(s, 'cross-white', px(0.885), py(0.30), 0.9)

add_note(s, "реж. заметка:\nбольше смеха, меньше повестки", px(0.05), py(0.045), -4)

card = add_card(s, px(0.055), py(0.16), Inches(6.6), Inches(5.55), rot=-1.1)
inner_l = px(0.055) + Inches(0.35)
inner_w = Inches(5.9)
y = py(0.16) + Inches(0.28)
add_text(s, inner_l, y, inner_w, Inches(0.4), "NEW NS TOURS", font='Caveat', size=20, bold=True,
         color=ACCENTS['rust']['light']); y += Inches(0.4)
add_text(s, inner_l, y, inner_w, Inches(1.95), "Грузия, где тимбилдинг вкуснее KPI",
         font='Yeseva One', size=32, color=INK, line_spacing=1.05); y += Inches(1.85)
add_text(s, inner_l, y, inner_w, Inches(0.7), "Пять однодневных форматов — выбираете тот, что подходит именно вашей команде.",
         font='Manrope', size=14, bold=True, color=RGBColor(0x3a, 0x27, 0x1f), line_spacing=1.15); y += Inches(0.68)
add_text(s, inner_l, y, inner_w, Inches(0.95), "Погружение в грузинскую культуру через гастрономию, ритуалы и живое взаимодействие. Каждый формат — законченный день, группа от 10 человек.",
         font='Manrope', size=12, color=RGBColor(0x40,0x35,0x2c), line_spacing=1.3); y += Inches(0.9)
add_text(s, inner_l, y, inner_w, Inches(0.6), "— да, это тот самый выезд, после которого в рабочем чате становится теплее.",
         font='Caveat', size=17, bold=True, color=ACCENTS['rust']['light'], line_spacing=1.05)

# ---------------------------------------------------------------
# SLIDE 2 — MENU
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_bg(s, 'mustard')
add_doodle(s, 'star', px(0.83), py(0.44), 2.1)
add_note(s, "все форматы — один день,\nгруппа от 10 человек", px(0.05), py(0.045), -3)

card = add_card(s, px(0.13), py(0.29), Inches(8.7), Inches(4.2), rot=-0.6)
inner_l = px(0.13) + Inches(0.4)
inner_w = Inches(7.9)
y = py(0.29) + Inches(0.34)
add_text(s, inner_l, y, inner_w, Inches(0.4), "Пять способов прочувствовать Грузию",
         font='Caveat', size=19, bold=True, color=ACCENTS['mustard']['light']); y += Inches(0.42)
add_text(s, inner_l, y, inner_w, Inches(0.7), "Выбираем формат", font='Yeseva One', size=30, color=INK); y += Inches(0.78)

rows = [
    ("01", "От лозы до бокала", "винный тимбилдинг с купажом", "$160/чел"),
    ("02", "5 чувств", "сенсорное погружение, специи, полифония", "$140/чел"),
    ("03", "Супра-баттл", "команда против тамады + кулинарный баттл", "$310/чел"),
    ("04", "Ретрит силы", "перезагрузка в Шато Агарани", "$175/чел"),
    ("05", "Утраченный рецепт", "квест по рынку Орбелиани", "$210/чел"),
]
row_h = Inches(0.5)
for num, name, hook, price in rows:
    add_text(s, inner_l, y, Inches(0.5), row_h, num, font='Yeseva One', size=14, color=ACCENTS['mustard']['light'], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, inner_l + Inches(0.55), y, Inches(2.1), row_h, name, font='Manrope', size=13.5, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, inner_l + Inches(2.75), y, Inches(3.9), row_h, hook, font='Manrope', size=11.5, color=RGBColor(0x45,0x3a,0x30), anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, inner_l + Inches(6.7), y, Inches(1.15), row_h, price, font='Manrope', size=13, bold=True, color=ACCENTS['mustard']['light'], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    line = s.shapes.add_connector(1, inner_l, y + row_h, inner_l + inner_w, y + row_h)
    line.line.color.rgb = RGBColor(0xC8, 0xB8, 0x98)
    line.line.width = Pt(0.75)
    y += row_h


# ---------------------------------------------------------------
# PROGRAM SLIDE TEMPLATE
# ---------------------------------------------------------------
def build_program_slide(num, accent, note, title, title_size, price, beats, humor,
                          photo_file, doodle_name, doodle_cx, doodle_cy, doodle_w,
                          card_rot, postcard_rot, note_rot, beats_size=12.5, h1_lines=2):
    s = prs.slides.add_slide(blank)
    add_bg(s, accent)
    add_doodle(s, doodle_name, doodle_cx, doodle_cy, doodle_w)
    add_note(s, note, px(0.05), py(0.04), note_rot)

    card_h = Inches(4.5 + 0.35 * max(0, h1_lines - 1) + 0.28 * max(0, len(beats) - 3))
    card = add_card(s, px(0.045), py(0.205), Inches(6.9), card_h, rot=card_rot)
    inner_l = px(0.045) + Inches(0.38)
    inner_w = Inches(6.1)
    y = py(0.205) + Inches(0.3)

    add_text(s, inner_l, y, inner_w, Inches(0.4), f"Программа {num}", font='Caveat', size=19,
             bold=True, color=ACCENTS[accent]['light']); y += Inches(0.4)

    h1_h = Inches(0.58 * h1_lines + 0.12)
    add_text(s, inner_l, y, inner_w, h1_h, title, font='Yeseva One', size=title_size, color=INK,
             line_spacing=1.05); y += h1_h + Inches(0.08)

    pill, ph = add_price_tag(s, inner_l, y, accent, price); y += ph + Inches(0.18)

    beats_tb = add_beats(s, inner_l, y, inner_w, beats, accent, size=beats_size)
    beats_est_h = Inches(0.30 * len(beats) + 0.06 * sum(1 for l, r in beats if len(l) + len(r) > 46))
    y += beats_est_h + Inches(0.12)

    add_text(s, inner_l, y, inner_w, Inches(0.55), "— " + humor, font='Caveat', size=16, bold=True,
             color=ACCENTS[accent]['light'], line_spacing=1.05)

    add_postcard(s, px(0.878), py(0.685), Inches(2.45), os.path.join(ASSETS, photo_file), postcard_rot)
    return s


build_program_slide(
    num="01", accent='grape', note="виноград давим ногами —\nтрадиция, не постановка",
    title="От лозы до бокала", title_size=32, h1_lines=1,
    price="$160",
    beats=[
        ("Сбор и давка винограда", "своими руками, по традиции"),
        ("Авторский купаж", "и презентация идеи эксперту"),
        ("Гастроужин", "игровая дегустация, угадать сорт в каждом блюде"),
        ("По желанию", "цифровой сертификат на своё вино"),
    ],
    humor="здесь спорят не о дедлайнах, а о нотах в аромате.",
    photo_file='ph_wine.jpg', doodle_name='grape', doodle_cx=px(0.895), doodle_cy=py(0.245), doodle_w=1.7,
    card_rot=0.9, postcard_rot=-4, note_rot=3,
)

build_program_slide(
    num="02", accent='teal', note="специи — с собой,\nв подарок",
    title="Грузия через 5 чувств", title_size=30, h1_lines=1,
    price="$140",
    beats=[
        ("Слух", "грузинская полифония вживую"),
        ("Запах и осязание", "мастер-класс по специям, набор с собой"),
        ("Вкус и зрение", "ужин с живым многоголосием"),
        ("Каждый уходит с", "авторским артефактом на память"),
    ],
    humor="план простой: вдохнули, попробовали, влюбились.",
    photo_file='ph_senses.jpg', doodle_name='senses', doodle_cx=px(0.895), doodle_cy=py(0.30), doodle_w=2.6,
    card_rot=1.1, postcard_rot=-3, note_rot=3,
)

build_program_slide(
    num="03", accent='rust', note="судят четыре шефа —\nвслепую",
    title="Супра-баттл: Корпоратив против Тамады", title_size=25, h1_lines=2,
    price="$310",
    beats=[
        ("Своя супра", "роли, тосты, атрибутика, сценарий вечера"),
        ("Kitchen Battle", "параллельно, под руководством четырёх шефов"),
        ("Слепая дегустация", "и награждение лучших"),
        ("С собой", "папаха, рог и крафтовая сумка с трофеями"),
    ],
    humor="единственный баттл, где красиво перебивать — часть жанра.",
    photo_file='ph_supra.jpg', doodle_name='horn', doodle_cx=px(0.90), doodle_cy=py(0.24), doodle_w=1.35,
    card_rot=0.9, postcard_rot=5, note_rot=3,
)

build_program_slide(
    num="04", accent='teal', note="Шато Агарани, Телави —\nна закате",
    title="Ретрит силы", title_size=32, h1_lines=1,
    price="$175",
    beats=[
        ("Дыхательные и телесные практики", "в полной тишине"),
        ("Место", "Шато Агарани, Телави, среди виноградников Кахетии"),
        ("Арт-элементы", "и полное отключение от рабочего чата"),
    ],
    humor="созвоны отменяются, глубокий вдох подтверждается.",
    photo_file='ph_retreat.jpg', doodle_name='sun', doodle_cx=px(0.905), doodle_cy=py(0.26), doodle_w=1.9,
    card_rot=1.3, postcard_rot=4, note_rot=3,
)

build_program_slide(
    num="05", accent='mustard', note="старт — рынок\nОрбелиани",
    title="В поисках утраченного рецепта", title_size=26, h1_lines=2,
    price="$210",
    beats=[
        ("Старт", "рынок Орбелиани"),
        ("Маршрут и задания", "по улицам Тбилиси, в поисках ингредиентов"),
        ("Финал", "кулинарный баттл, общее блюдо из находок"),
    ],
    humor="командировка, в которой все рады потеряться на рынке.",
    photo_file='ph_quest.jpg', doodle_name='route', doodle_cx=px(0.845), doodle_cy=py(0.235), doodle_w=2.3,
    card_rot=1.5, postcard_rot=5, note_rot=3,
)

# ---------------------------------------------------------------
# SLIDE 8 — CONTACT
# ---------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_bg(s, 'rust')
add_doodle(s, 'wreath', px(0.885), py(0.28), 2.5)
add_note(s, "выбирайте формат —\nостальное решим мы", px(0.05), py(0.04), -3)

card = add_card(s, px(0.045), py(0.28), Inches(6.6), Inches(4.35), rot=-1.2)
inner_l = px(0.045) + Inches(0.38)
inner_w = Inches(5.85)
y = py(0.28) + Inches(0.3)
add_text(s, inner_l, y, inner_w, Inches(0.4), "Финал", font='Caveat', size=19, bold=True,
         color=ACCENTS['rust']['light']); y += Inches(0.42)
add_text(s, inner_l, y, inner_w, Inches(0.75), "Выбираем формат?", font='Yeseva One', size=34, color=INK); y += Inches(0.78)
add_text(s, inner_l, y, inner_w, Inches(0.45), "Посчитаем точную смету под вашу группу и даты.",
         font='Manrope', size=14, bold=True, color=RGBColor(0x3a,0x27,0x1f)); y += Inches(0.45)
add_text(s, inner_l, y, inner_w, Inches(0.7), "Пять форматов, один день, группа от 10 человек. Уезжаете с сувенирами, возвращаетесь с историями.",
         font='Manrope', size=12, color=RGBColor(0x40,0x35,0x2c), line_spacing=1.3); y += Inches(0.75)

line = s.shapes.add_connector(1, inner_l, y, inner_l + inner_w, y)
line.line.color.rgb = RGBColor(0xC8, 0xB8, 0x98); line.line.width = Pt(0.75)
y += Inches(0.16)

add_text(s, inner_l, y, inner_w, Inches(0.4), "Нини Салуквадзе", font='Yeseva One', size=19, color=INK); y += Inches(0.42)
add_text(s, inner_l, y, inner_w, Inches(0.3), "Генеральный директор", font='Manrope', size=12,
         color=RGBColor(0x55,0x48,0x3c)); y += Inches(0.32)
add_text(s, inner_l, y, inner_w, Inches(0.35), "+995 551 91 50 90  ·  info@nstours.ge  ·  nstours.ge",
         font='Manrope', size=13, bold=True, color=ACCENTS['rust']['light'])

add_postcard(s, px(0.875), py(0.685), Inches(2.45), os.path.join(ASSETS, 'ph_contact.jpg'), -4)

out_dir = os.path.join(BASE, 'build')
os.makedirs(out_dir, exist_ok=True)
pptx_out = os.path.join(out_dir, 'NS-Tours-Concept2.pptx')
prs.save(pptx_out)
print("Saved full 8-slide deck:", pptx_out)
